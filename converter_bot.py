# ============================================================
# EDUBOT v12 — Professional Ta'lim Boti
# 3 tilda (O'zbek, Rus, Ingliz) | Imlo xatosiz | Barcha xizmatlar
# ============================================================
import telebot, os, sqlite3, logging, tempfile, shutil, requests, re, math
from telebot import types
from datetime import datetime
from io import BytesIO

# ============================================================
# SOZLAMALAR
# ============================================================
BOT_TOKEN      = os.environ.get("BOT_TOKEN",      "")
ADMIN_ID       = int(os.environ.get("ADMIN_ID",   "1113404703"))
ADMIN_USERNAME = os.environ.get("ADMIN_USERNAME", "@abdurakhmon02")
DONATE_URL     = os.environ.get("DONATE_URL",     "https://click.uz")
DONATE_CARD    = os.environ.get("DONATE_CARD",    "9860 0609 2665 0809")
DONATE_CLICK   = os.environ.get("DONATE_CLICK",   "+998 94 975 03 04")
CLAUDE_API_KEY = os.environ.get("CLAUDE_API_KEY", "")
UNSPLASH_KEY   = os.environ.get("UNSPLASH_KEY",   "")
SONNET_MODEL   = os.environ.get("SONNET_MODEL", "claude-sonnet-4-6")
HAIKU_MODEL    = os.environ.get("HAIKU_MODEL", "claude-haiku-4-5-20251001")

def gp(n, d): return int(os.environ.get(n, str(d)))
PRICE_PAGE     = gp("PRICE_PAGE",    300)
PRICE_KURS     = gp("PRICE_KURS",    250)
PRICE_MUSTAQIL = gp("PRICE_MUSTAQIL",280)
PRICE_MAQOLA   = gp("PRICE_MAQOLA",  400)
PRICE_SLIDE    = gp("PRICE_SLIDE",   300)
PRICE_TEST     = gp("PRICE_TEST",    150)
BONUS_FIRST    = gp("BONUS_FIRST",  3000)


# ============================================================
# MAJBURIY OBUNA TIZIMI
# ============================================================
SUB_ENABLED = os.environ.get("SUB_ENABLED", "0") == "1"

def get_sub_channels():
    """DB dan obuna kanallarini olish"""
    try:
        c = sqlite3.connect("edubot.db"); cur = c.cursor()
        cur.execute("SELECT channel_id, channel_name FROM sub_channels WHERE active=1")
        rows = cur.fetchall(); c.close()
        return rows
    except: return []

def add_sub_channel(channel_id, channel_name):
    try:
        c = sqlite3.connect("edubot.db"); cur = c.cursor()
        cur.execute("INSERT OR REPLACE INTO sub_channels(channel_id, channel_name, active) VALUES(?,?,1)",
            (channel_id, channel_name))
        c.commit(); c.close(); return True
    except: return False

def remove_sub_channel(channel_id):
    try:
        c = sqlite3.connect("edubot.db"); cur = c.cursor()
        cur.execute("DELETE FROM sub_channels WHERE channel_id=?", (channel_id,))
        c.commit(); c.close(); return True
    except: return False

def check_subscription(uid):
    """Foydalanuvchi barcha kanallarga obuna bo'lganmi"""
    if not SUB_ENABLED: return True
    channels = get_sub_channels()
    if not channels: return True
    for ch_id, ch_name in channels:
        try:
            member = bot.get_chat_member(ch_id, uid)
            if member.status in ['left', 'kicked', 'banned']:
                return False
        except: return False
    return True

def sub_check_kb():
    """Obuna tekshirish tugmasi"""
    channels = get_sub_channels()
    kb = types.InlineKeyboardMarkup(row_width=1)
    for ch_id, ch_name in channels:
        try:
            invite = bot.export_chat_invite_link(ch_id)
        except:
            invite = f"https://t.me/{str(ch_id).lstrip('-100')}"
        kb.add(types.InlineKeyboardButton(f"📢 {ch_name}", url=invite))
    kb.add(types.InlineKeyboardButton("✅ Obunani tekshirish", callback_data="check_sub"))
    return kb

def require_sub(uid, func, *args, **kwargs):
    """Obunani tekshirib, agar obuna bo'lmasa xabar yuborish"""
    if check_subscription(uid):
        return func(*args, **kwargs)
    channels = get_sub_channels()
    ch_list = "\n".join([f"• {name}" for _, name in channels])
    bot.send_message(uid,
        f"⚠️ *Botdan foydalanish uchun quyidagi kanallarga obuna bo'ling:*\n\n{ch_list}\n\n"
        f"Obuna bo'lgach ✅ tugmasini bosing.",
        parse_mode="Markdown",
        reply_markup=sub_check_kb())
    return None

logging.basicConfig(format="%(asctime)s - %(levelname)s - %(message)s", level=logging.INFO)
logger = logging.getLogger(__name__)
bot = telebot.TeleBot(BOT_TOKEN)

# ============================================================
# KO'P TILLI LUGAT
# ============================================================
TEXTS = {
    "uz": {
        "welcome": "👋 Xush kelibsiz, {name}!\n\n🎓 *EduBot* — Ta'lim yordamchingiz!\n\n📚 Xizmatlar:\n📄 Referat | 📝 Kurs ishi | 📋 Mustaqil ish\n📰 Maqola | 📊 Prezentatsiya | ✅ Test\n✏️ Imlo tuzatish | 🔄 Konvertatsiya",
        "bonus": "\n\n🎁 *{amount:,} so'm bonus berildi!*",
        "choose_lang": "\n\nTilni tanlang:",
        "menu": "📋 Asosiy menyu:",
        "enter_topic": "📝 Mavzuni kiriting:",
        "enter_pages": "📄 Necha bet? (5-100)\n💰 1 bet = {price:,} so'm",
        "enter_slides": "🎯 Necha slayd? (10-50)\n💰 1 slayd = {price:,} so'm",
        "enter_count": "🔢 Nechta savol? (10-1000)\n💰 1 savol = {price:,} so'm",
        "ask_name": "👤 Ism va familiyangizni kiriting:",
        "ask_univ": "🏛 Universitetingiz nomi:",
        "ask_faculty": "📚 Fakultetingiz:",
        "ask_year": "📅 Nechinchi kurs?",
        "ask_teacher": "👩‍🏫 O'qituvchi ismi:",
        "ask_subject": "📖 Fan nomi:",
        "ask_city": "🏙 Shahar:",
        "optional": "(ixtiyoriy — o'tkazib yuborish mumkin)",
        "ask_lang": "🌐 Qaysi tilda?",
        "ask_format": "📁 Format tanlang:",
        "ask_plans": "📋 Nechta bo'lim (reja) bo'lsin?",
        "ask_template": "🎨 Shablon tanlang:",
        "preparing": "⏳ Tayyorlanmoqda... Iltimos kuting.",
        "ready": "✅ Tayyor!\n💰 Balans: {bal:,} so'm",
        "error": "❌ Xatolik yuz berdi. Pul qaytarildi.",
        "no_balance": "❌ Mablag' yetarli emas!\nKerakli: {need:,} so'm\nBalans: {bal:,} so'm\n\nBalansni to'ldiring:",
        "balance_info": "💰 *Balansingiz: {bal:,} so'm*",
        "topup": "💳 Balans to'ldirish",
        "back": "◀️ Orqaga",
        "home": "🏠 Menyu",
        "skip": "⏭ O'tkazib yuborish",
        "lang_set": "✅ Til o'rnatildi: O'zbek 🇺🇿",
        "no_orders": "📦 Hozircha buyurtmalaringiz yo'q.\n\nBirinchi buyurtmangizni bering!",
        "orders_title": "📦 *Oxirgi buyurtmalaringiz:*\n\n",
        "write_topic": "📝 Mavzuni kiriting:",
        "img_accept": "🖼 Rasm qabul qilindi! Qaysi slaydga qo'yish kerak?",
        "img_count": "✅ {n} ta rasm qabul qilindi. /done yozing.",
        "converting": "⏳ Konvertatsiya amalga oshirilmoqda...",
        "wrong_format": "❌ Noto'g'ri format.",
        "imlo_check": "⏳ Imlo tekshirilmoqda...",
        "imlo_done": "✅ *Imlo tuzatildi:*\n\n",
        "imlo_notfound": "❌ Matn topilmadi.",
        "imlo_prompt": "✏️ Matnni yozing yoki fayl (PDF/TXT) yuboring:",
        "ai_img": "🤖 AI avtomatik rasm qo'ysin",
        "user_img": "🖼 O'zim rasm yuklаyman",
        "no_img": "❌ Rasmsiz davom etish",
        "help_text": "❓ *Yordam*\n\n📄 /referat — Referat\n📝 /kursishi — Kurs ishi\n📋 /mustaqilish — Mustaqil ish\n📰 /maqola — Maqola\n📊 /prezentatsiya — Prezentatsiya\n✅ /test — Test\n✏️ /imlo — Imlo tuzatish\n🔄 /konvertatsiya — Konvertatsiya\n💰 /balans — Balans",
        "prices": "💵 *Narxlar:*\n📄 Referat: {p1:,} so'm/bet\n📝 Kurs ishi: {p2:,} so'm/bet\n📋 Mustaqil: {p3:,} so'm/bet\n📰 Maqola: {p4:,} so'm/bet\n📊 Prezentatsiya: {p5:,} so'm/slayd\n✅ Test: {p6:,} so'm/savol",
        "btn_referat": "📄 Referat",
        "btn_kurs": "📝 Kurs ishi",
        "btn_mustaqil": "📋 Mustaqil ish",
        "btn_maqola": "📰 Maqola",
        "btn_prez": "📊 Prezentatsiya",
        "btn_test": "✅ Test",
        "btn_imlo": "✏️ Imlo tuzatish",
        "btn_konv": "🔄 Konvertatsiya",
        "btn_balans": "💰 Balans",
        "btn_orders": "📦 Buyurtmalarim",
        "btn_referral": "👥 Referal",
        "btn_donat": "💝 Donat",
        "btn_help": "❓ Yordam",
        "btn_admin": "👨‍💼 Admin",
        "slide_pages": "slayd",
        "bet_pages": "bet",
        "savol_pages": "savol",
    },
    "ru": {
        "welcome": "👋 Добро пожаловать, {name}!\n\n🎓 *EduBot* — Ваш учебный помощник!\n\n📚 Услуги:\n📄 Реферат | 📝 Курсовая | 📋 Самостоятельная\n📰 Статья | 📊 Презентация | ✅ Тест\n✏️ Проверка орфографии | 🔄 Конвертация",
        "bonus": "\n\n🎁 *Начислено {amount:,} сум бонус!*",
        "choose_lang": "\n\nВыберите язык:",
        "menu": "📋 Главное меню:",
        "enter_topic": "📝 Введите тему:",
        "enter_pages": "📄 Сколько страниц? (5-100)\n💰 1 стр = {price:,} сум",
        "enter_slides": "🎯 Сколько слайдов? (10-50)\n💰 1 слайд = {price:,} сум",
        "enter_count": "🔢 Сколько вопросов? (10-1000)\n💰 1 вопрос = {price:,} сум",
        "ask_name": "👤 Введите имя и фамилию:",
        "ask_univ": "🏛 Название университета:",
        "ask_faculty": "📚 Факультет:",
        "ask_year": "📅 Курс (1-4)?",
        "ask_teacher": "👩‍🏫 Имя преподавателя:",
        "ask_subject": "📖 Название предмета:",
        "ask_city": "🏙 Город:",
        "optional": "(необязательно — можно пропустить)",
        "ask_lang": "🌐 На каком языке?",
        "ask_format": "📁 Выберите формат:",
        "ask_plans": "📋 Сколько разделов?",
        "ask_template": "🎨 Выберите шаблон:",
        "preparing": "⏳ Подготавливается... Пожалуйста, подождите.",
        "ready": "✅ Готово!\n💰 Баланс: {bal:,} сум",
        "error": "❌ Произошла ошибка. Деньги возвращены.",
        "no_balance": "❌ Недостаточно средств!\nНужно: {need:,} сум\nБаланс: {bal:,} сум\n\nПополните баланс:",
        "balance_info": "💰 *Ваш баланс: {bal:,} сум*",
        "topup": "💳 Пополнить баланс",
        "back": "◀️ Назад",
        "home": "🏠 Меню",
        "skip": "⏭ Пропустить",
        "lang_set": "✅ Язык установлен: Русский 🇷🇺",
        "no_orders": "📦 У вас пока нет заказов.\n\nСделайте первый заказ!",
        "orders_title": "📦 *Ваши последние заказы:*\n\n",
        "write_topic": "📝 Введите тему:",
        "img_accept": "🖼 Фото принято! На какой слайд добавить?",
        "img_count": "✅ Принято {n} фото. Напишите /done.",
        "converting": "⏳ Конвертация выполняется...",
        "wrong_format": "❌ Неверный формат.",
        "imlo_check": "⏳ Проверка орфографии...",
        "imlo_done": "✅ *Орфография исправлена:*\n\n",
        "imlo_notfound": "❌ Текст не найден.",
        "imlo_prompt": "✏️ Напишите текст или отправьте файл (PDF/TXT):",
        "ai_img": "🤖 AI добавит фото автоматически",
        "user_img": "🖼 Загружу сам",
        "no_img": "❌ Без фото",
        "help_text": "❓ *Помощь*\n\n📄 /referat — Реферат\n📝 /kursishi — Курсовая\n📋 /mustaqilish — Самостоятельная\n📰 /maqola — Статья\n📊 /prezentatsiya — Презентация\n✅ /test — Тест\n✏️ /imlo — Орфография\n🔄 /konvertatsiya — Конвертация\n💰 /balans — Баланс",
        "prices": "💵 *Цены:*\n📄 Реферат: {p1:,}/стр\n📝 Курсовая: {p2:,}/стр\n📋 Самост.: {p3:,}/стр\n📰 Статья: {p4:,}/стр\n📊 Презент.: {p5:,}/слайд\n✅ Тест: {p6:,}/вопрос",
        "btn_referat": "📄 Реферат",
        "btn_kurs": "📝 Курсовая",
        "btn_mustaqil": "📋 Самостоятельная",
        "btn_maqola": "📰 Статья",
        "btn_prez": "📊 Презентация",
        "btn_test": "✅ Тест",
        "btn_imlo": "✏️ Орфография",
        "btn_konv": "🔄 Конвертация",
        "btn_balans": "💰 Баланс",
        "btn_orders": "📦 Мои заказы",
        "btn_referral": "👥 Реферал",
        "btn_donat": "💝 Донат",
        "btn_help": "❓ Помощь",
        "btn_admin": "👨‍💼 Админ",
        "slide_pages": "слайдов",
        "bet_pages": "стр",
        "savol_pages": "вопросов",
    },
    "en": {
        "welcome": "👋 Welcome, {name}!\n\n🎓 *EduBot* — Your Educational Assistant!\n\n📚 Services:\n📄 Essay | 📝 Course Work | 📋 Independent Work\n📰 Article | 📊 Presentation | ✅ Test\n✏️ Spell Check | 🔄 Convert",
        "bonus": "\n\n🎁 *{amount:,} sum bonus added!*",
        "choose_lang": "\n\nChoose language:",
        "menu": "📋 Main menu:",
        "enter_topic": "📝 Enter topic:",
        "enter_pages": "📄 How many pages? (5-100)\n💰 1 page = {price:,} sum",
        "enter_slides": "🎯 How many slides? (10-50)\n💰 1 slide = {price:,} sum",
        "enter_count": "🔢 How many questions? (10-1000)\n💰 1 question = {price:,} sum",
        "ask_name": "👤 Enter your name:",
        "ask_univ": "🏛 University name:",
        "ask_faculty": "📚 Faculty:",
        "ask_year": "📅 Year (1-4)?",
        "ask_teacher": "👩‍🏫 Teacher's name:",
        "ask_subject": "📖 Subject name:",
        "ask_city": "🏙 City:",
        "optional": "(optional — you can skip)",
        "ask_lang": "🌐 In which language?",
        "ask_format": "📁 Choose format:",
        "ask_plans": "📋 How many sections?",
        "ask_template": "🎨 Choose template:",
        "preparing": "⏳ Preparing... Please wait.",
        "ready": "✅ Done!\n💰 Balance: {bal:,} sum",
        "error": "❌ An error occurred. Money refunded.",
        "no_balance": "❌ Insufficient balance!\nNeeded: {need:,} sum\nBalance: {bal:,} sum\n\nTop up balance:",
        "balance_info": "💰 *Your balance: {bal:,} sum*",
        "topup": "💳 Top up balance",
        "back": "◀️ Back",
        "home": "🏠 Menu",
        "skip": "⏭ Skip",
        "lang_set": "✅ Language set: English 🇬🇧",
        "no_orders": "📦 No orders yet.\n\nMake your first order!",
        "orders_title": "📦 *Your recent orders:*\n\n",
        "write_topic": "📝 Enter topic:",
        "img_accept": "🖼 Photo accepted! Which slide to add to?",
        "img_count": "✅ {n} photos received. Type /done.",
        "converting": "⏳ Converting...",
        "wrong_format": "❌ Wrong format.",
        "imlo_check": "⏳ Checking spelling...",
        "imlo_done": "✅ *Spelling corrected:*\n\n",
        "imlo_notfound": "❌ Text not found.",
        "imlo_prompt": "✏️ Write text or send file (PDF/TXT):",
        "ai_img": "🤖 AI auto add photos",
        "user_img": "🖼 I'll upload myself",
        "no_img": "❌ Without photos",
        "help_text": "❓ *Help*\n\n📄 /referat — Essay\n📝 /kursishi — Course Work\n📋 /mustaqilish — Independent\n📰 /maqola — Article\n📊 /prezentatsiya — Presentation\n✅ /test — Test\n✏️ /imlo — Spell Check\n🔄 /konvertatsiya — Convert\n💰 /balans — Balance",
        "prices": "💵 *Prices:*\n📄 Essay: {p1:,}/page\n📝 Course: {p2:,}/page\n📋 Independent: {p3:,}/page\n📰 Article: {p4:,}/page\n📊 Presentation: {p5:,}/slide\n✅ Test: {p6:,}/question",
        "btn_referat": "📄 Essay",
        "btn_kurs": "📝 Course Work",
        "btn_mustaqil": "📋 Independent",
        "btn_maqola": "📰 Article",
        "btn_prez": "📊 Presentation",
        "btn_test": "✅ Test",
        "btn_imlo": "✏️ Spell Check",
        "btn_konv": "🔄 Convert",
        "btn_balans": "💰 Balance",
        "btn_orders": "📦 My Orders",
        "btn_referral": "👥 Referral",
        "btn_donat": "💝 Donate",
        "btn_help": "❓ Help",
        "btn_admin": "👨‍💼 Admin",
        "slide_pages": "slides",
        "bet_pages": "pages",
        "savol_pages": "questions",
    }
}

def t(uid, key, **kwargs):
    """Foydalanuvchi tiliga mos matn"""
    lang = get_lang(uid)
    text = TEXTS.get(lang, TEXTS["uz"]).get(key, TEXTS["uz"].get(key, key))
    if kwargs:
        try: text = text.format(**kwargs)
        except: pass
    return text

LN = {"uz": "o'zbek", "ru": "rus", "en": "ingliz"}

# ============================================================
# DATABASE
# ============================================================
def init_db():
    c = sqlite3.connect("edubot.db")
    cur = c.cursor()
    cur.execute("""CREATE TABLE IF NOT EXISTS users(
        telegram_id INTEGER PRIMARY KEY, username TEXT, first_name TEXT,
        lang TEXT DEFAULT 'uz', balance INTEGER DEFAULT 0, joined_at TEXT)""")
    cur.execute("""CREATE TABLE IF NOT EXISTS orders(
        telegram_id INTEGER PRIMARY KEY, state TEXT, data TEXT, updated_at TEXT)""")
    cur.execute("""CREATE TABLE IF NOT EXISTS stats(
        id INTEGER PRIMARY KEY AUTOINCREMENT, telegram_id INTEGER,
        action TEXT, detail TEXT, income INTEGER DEFAULT 0, created_at TEXT)""")
    cur.execute("""CREATE TABLE IF NOT EXISTS buyurtmalar(
        id INTEGER PRIMARY KEY AUTOINCREMENT, telegram_id INTEGER,
        tur TEXT, mavzu TEXT, format TEXT, sahifalar TEXT,
        narx INTEGER, status TEXT DEFAULT 'done', order_data TEXT, created_at TEXT)""")
    # Eski jadvalga ustun qo'shish (agar yo'q bo'lsa)
    try:
        cur.execute("ALTER TABLE buyurtmalar ADD COLUMN status TEXT DEFAULT 'done'")
    except: pass
    try:
        cur.execute("ALTER TABLE buyurtmalar ADD COLUMN order_data TEXT")
    except: pass
    cur.execute("""CREATE TABLE IF NOT EXISTS sub_channels(
        channel_id TEXT PRIMARY KEY, channel_name TEXT, active INTEGER DEFAULT 1)""")
    cur.execute("""CREATE TABLE IF NOT EXISTS referrals(
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        referrer_id INTEGER,
        referred_id INTEGER UNIQUE,
        bonus_paid INTEGER DEFAULT 0,
        created_at TEXT)""")
    cur.execute("""CREATE TABLE IF NOT EXISTS topup_requests(
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        telegram_id INTEGER,
        amount INTEGER,
        status TEXT DEFAULT 'pending',
        admin_note TEXT,
        created_at TEXT)""")
    try:
        cur.execute("ALTER TABLE users ADD COLUMN referred_by INTEGER DEFAULT NULL")
    except: pass
    c.commit(); c.close()

import json

def save_order(uid, state, data):
    try:
        c = sqlite3.connect("edubot.db"); cur = c.cursor()
        cur.execute("INSERT OR REPLACE INTO orders VALUES(?,?,?,?)",
            (uid, state, json.dumps(data, ensure_ascii=False),
             datetime.now().strftime("%d.%m.%Y %H:%M")))
        c.commit(); c.close()
    except: pass

def load_order(uid):
    try:
        c = sqlite3.connect("edubot.db"); cur = c.cursor()
        cur.execute("SELECT state,data FROM orders WHERE telegram_id=?", (uid,))
        row = cur.fetchone(); c.close()
        if row: return row[0], json.loads(row[1])
    except: pass
    return None, {}

def clear_order(uid):
    try:
        c = sqlite3.connect("edubot.db"); cur = c.cursor()
        cur.execute("DELETE FROM orders WHERE telegram_id=?", (uid,))
        c.commit(); c.close()
    except: pass

def get_user(uid):
    c = sqlite3.connect("edubot.db"); cur = c.cursor()
    cur.execute("SELECT * FROM users WHERE telegram_id=?", (uid,))
    row = cur.fetchone(); c.close()
    return {"id": row[0], "username": row[1], "first_name": row[2],
            "lang": row[3], "balance": row[4]} if row else None

def reg_user(uid, uname, fname, lang="uz"):
    c = sqlite3.connect("edubot.db"); cur = c.cursor()
    cur.execute("SELECT telegram_id FROM users WHERE telegram_id=?", (uid,))
    ex = cur.fetchone()
    if not ex:
        cur.execute("INSERT INTO users VALUES(?,?,?,?,?,?,?)",
            (uid, uname, fname, lang, BONUS_FIRST,
             datetime.now().strftime("%d.%m.%Y %H:%M"), None))
        c.commit(); c.close(); return True
    cur.execute("UPDATE users SET username=?,first_name=? WHERE telegram_id=?",
        (uname, fname, uid))
    c.commit(); c.close(); return False

def get_lang(uid):
    u = get_user(uid); return u["lang"] if u else "uz"

def set_lang(uid, lang):
    c = sqlite3.connect("edubot.db"); cur = c.cursor()
    cur.execute("UPDATE users SET lang=? WHERE telegram_id=?", (lang, uid))
    c.commit(); c.close()

def get_balance(uid):
    u = get_user(uid); return u["balance"] if u else 0

def deduct(uid, amt):
    c = sqlite3.connect("edubot.db"); cur = c.cursor()
    cur.execute("UPDATE users SET balance=balance-? WHERE telegram_id=?", (amt, uid))
    c.commit(); c.close()

def add_bal(uid, amt):
    c = sqlite3.connect("edubot.db"); cur = c.cursor()
    cur.execute("UPDATE users SET balance=balance+? WHERE telegram_id=?", (amt, uid))
    c.commit(); c.close()

def log_act(uid, action, detail="", income=0):
    c = sqlite3.connect("edubot.db"); cur = c.cursor()
    cur.execute("INSERT INTO stats(telegram_id,action,detail,income,created_at) VALUES(?,?,?,?,?)",
        (uid, action, detail, income, datetime.now().strftime("%d.%m.%Y %H:%M")))
    c.commit(); c.close()

def save_buyurtma(uid, tur, mavzu, fmt, sah, narx, status="done", order_data=None):
    try:
        c = sqlite3.connect("edubot.db"); cur = c.cursor()
        # Bu foydalanuvchining birinchi buyurtmasimi?
        cur.execute("SELECT COUNT(*) FROM buyurtmalar WHERE telegram_id=? AND status='done'", (uid,))
        is_first = cur.fetchone()[0] == 0
        cur.execute("""INSERT INTO buyurtmalar(telegram_id,tur,mavzu,format,sahifalar,narx,status,order_data,created_at) 
            VALUES(?,?,?,?,?,?,?,?,?)""",
            (uid, tur, mavzu, fmt, str(sah), narx, status,
             json.dumps(order_data or {}, ensure_ascii=False),
             datetime.now().strftime("%d.%m.%Y %H:%M")))
        c.commit()
        row_id = cur.lastrowid
        c.close()
        # Birinchi buyurtmada referal bonus
        if is_first and status == "done":
            import threading
            threading.Thread(target=pay_referral_bonus, args=(uid,), daemon=True).start()
        return row_id
    except Exception as e:
        logger.error(f"save_buyurtma: {e}"); return None


def save_pending_and_notify(uid, svc, topic, fmt, amount, total, ud):
    """Pul yetmaganda buyurtmani saqlash va xabar yuborish"""
    order_data = dict(ud)
    order_data["svc"] = svc
    save_pending_buyurtma(uid, svc, topic, fmt, amount, total, order_data)
    bal = get_balance(uid)
    kb2 = types.InlineKeyboardMarkup(row_width=1)
    kb2.add(types.InlineKeyboardButton(t(uid, "topup"), callback_data="topup"))
    kb2.add(types.InlineKeyboardButton("📋 Buyurtmalarim", callback_data="my_orders"))
    svc_names = {"referat":"📄 Referat","kurs":"📝 Kurs ishi","mustaqil":"📋 Mustaqil ish",
                 "maqola":"📰 Maqola","prez":"📊 Prezentatsiya","test":"✅ Test"}
    sl = "slayd" if svc == "prez" else ("savol" if svc == "test" else "bet")
    bot.send_message(uid,
        f"💰 *Hisobingizda mablag\' yetarli emas!*\n\n"
        f"📌 {svc_names.get(svc, svc)}: *{topic}*\n"
        f"📊 Hajm: {amount} {sl}\n"
        f"💵 Kerak: *{total:,} so\'m*\n"
        f"💳 Balans: *{bal:,} so\'m*\n\n"
        f"✅ Buyurtma *Buyurtmalarim* bo\'limiga saqlandi!\n"
        f"Hisobni to\'ldirgach, u yerdan davom eting.",
        parse_mode="Markdown", reply_markup=kb2)
    cst(uid)

def save_pending_buyurtma(uid, tur, mavzu, fmt, sah, narx, order_data):
    """Yakunlanmagan buyurtmani saqlash"""
    try:
        c = sqlite3.connect("edubot.db"); cur = c.cursor()
        # Avval eski pending buyurtmani o'chirish
        cur.execute("DELETE FROM buyurtmalar WHERE telegram_id=? AND status='pending' AND tur=?", (uid, tur))
        cur.execute("""INSERT INTO buyurtmalar(telegram_id,tur,mavzu,format,sahifalar,narx,status,order_data,created_at) 
            VALUES(?,?,?,?,?,?,?,?,?)""",
            (uid, tur, mavzu, fmt, str(sah), narx, "pending",
             json.dumps(order_data, ensure_ascii=False),
             datetime.now().strftime("%d.%m.%Y %H:%M")))
        c.commit()
        row_id = cur.lastrowid
        c.close()
        return row_id
    except Exception as e:
        logger.error(f"save_pending: {e}")
        return None

def get_pending_buyurtmalar(uid):
    """Yakunlanmagan buyurtmalarni olish"""
    try:
        c = sqlite3.connect("edubot.db"); cur = c.cursor()
        cur.execute("""SELECT id,tur,mavzu,format,sahifalar,narx,order_data,created_at 
            FROM buyurtmalar WHERE telegram_id=? AND status='pending' 
            ORDER BY id DESC LIMIT 5""", (uid,))
        rows = cur.fetchall(); c.close()
        return rows
    except: return []

def complete_buyurtma(order_id):
    """Buyurtmani yakunlash"""
    try:
        c = sqlite3.connect("edubot.db"); cur = c.cursor()
        cur.execute("UPDATE buyurtmalar SET status='done' WHERE id=?", (order_id,))
        c.commit(); c.close()
    except: pass

def delete_buyurtma(order_id):
    """Buyurtmani o'chirish"""
    try:
        c = sqlite3.connect("edubot.db"); cur = c.cursor()
        cur.execute("DELETE FROM buyurtmalar WHERE id=?", (order_id,))
        c.commit(); c.close()
    except: pass

def get_buyurtmalar(uid):
    try:
        c = sqlite3.connect("edubot.db"); cur = c.cursor()
        cur.execute("""SELECT id,tur,mavzu,format,sahifalar,narx,status,created_at 
            FROM buyurtmalar WHERE telegram_id=? 
            ORDER BY CASE WHEN status='pending' THEN 0 ELSE 1 END, id DESC LIMIT 15""", (uid,))
        rows = cur.fetchall(); c.close(); return rows
    except: return []

def all_users():
    c = sqlite3.connect("edubot.db"); cur = c.cursor()
    cur.execute("SELECT telegram_id FROM users")
    rows = cur.fetchall(); c.close()
    return [r[0] for r in rows]

def get_stats():
    c = sqlite3.connect("edubot.db"); cur = c.cursor()
    cur.execute("SELECT COUNT(*) FROM users"); u = cur.fetchone()[0]
    cur.execute("SELECT COUNT(*) FROM stats WHERE action IN ('referat','kurs','mustaqil','maqola','prez','test')"); w = cur.fetchone()[0]
    cur.execute("SELECT COUNT(*) FROM stats WHERE action='conv'"); cv = cur.fetchone()[0]
    cur.execute("SELECT COALESCE(SUM(income),0) FROM stats"); i = cur.fetchone()[0]
    c.close(); return u, w, cv, i

# ============================================================
# HOLAT BOSHQARUVI
# ============================================================
ST, UD, UI = {}, {}, {}
HISTORY = {}

def sst(uid, s, **kw):
    prev = ST.get(uid)
    if prev and prev != s:
        HISTORY.setdefault(uid, []).append(prev)
        if len(HISTORY[uid]) > 20:
            HISTORY[uid] = HISTORY[uid][-20:]
    ST[uid] = s
    UD.setdefault(uid, {}).update(kw)
    save_order(uid, s, UD.get(uid, {}))

def gst(uid): return ST.get(uid)

def cst(uid):
    ST.pop(uid, None)
    HISTORY.pop(uid, None)
    clear_order(uid)

def go_back(uid):
    h = HISTORY.get(uid, [])
    if not h: return None
    prev = h.pop()
    HISTORY[uid] = h
    ST[uid] = prev
    return prev

def restore_state(uid):
    state, data = load_order(uid)
    if state and data:
        ST[uid] = state; UD[uid] = data; return True
    return False

def build_info(ud):
    parts = []
    for k, lbl in [("full_name","Muallif"), ("subject","Fan"), ("university","Universitet"),
                   ("faculty","Fakultet"), ("year","Kurs"), ("teacher","O'qituvchi"), ("city","Shahar")]:
        if ud.get(k): parts.append(f"{lbl}: {ud[k]}")
    return "\n".join(parts)


# ============================================================
# DIAGRAMMA MODULI
# ============================================================
# ============================================================
# DIAGRAMMA VA INFOGRAFIKA MODULI
# Mavzuga qarab avtomatik diagramma qo'yadi
# ============================================================
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
import re, json

def rgb(r,g,b): return RGBColor(r,g,b)

# ============================================================
# DIAGRAMMA KERAKMI TEKSHIRISH
# ============================================================
DIAGRAM_KEYWORDS = {
    "bar_chart": [
        "iqtisod","statistik","o'sish","pasayish","daromad","ishlab chiqarish",
        "eksport","import","yalpi","GDP","BYM","foiz","%","million","milliard",
        "экономика","рост","снижение","доход","производство","экспорт",
        "economic","growth","decline","revenue","production","export","import",
        "miqdor","hajm","ko'rsatkich","natija","taqqoslash"
    ],
    "pie_chart": [
        "ulush","taqsimot","tarkib","qism","struktura","nisbat",
        "доля","распределение","структура","часть",
        "share","distribution","structure","portion","percentage",
        "foiz ulushi","qancha qism","nechta"
    ],
    "line_chart": [
        "dinamika","trend","vaqt","yil","oy","davr","o'zgarish","rivojlanish",
        "динамика","тренд","время","год","месяц","период","изменение",
        "timeline","trend","period","change","development","history",
        "tarix","o'tgan","kelajak","prognoz"
    ],
    "infographic": [
        "bosqich","jarayon","ketma-ket","qadamlar","sxema","algoritm",
        "этап","процесс","шаг","схема","алгоритм",
        "step","process","stage","flow","scheme","algorithm",
        "anatomiya","tuzilish","qism","element"
    ]
}

def detect_diagram_type(topic, slide_title, slide_text):
    """Mavzu va matnga qarab diagramma turini aniqlash"""
    combined = f"{topic} {slide_title} {slide_text}".lower()
    
    scores = {"bar_chart": 0, "pie_chart": 0, "line_chart": 0, "infographic": 0, "none": 0}
    
    for dtype, keywords in DIAGRAM_KEYWORDS.items():
        for kw in keywords:
            if kw.lower() in combined:
                scores[dtype] += 1
    
    max_type = max(scores, key=scores.get)
    max_score = scores[max_type]
    
    if max_score == 0:
        return "none"
    return max_type

def parse_diagram_data(claude_text):
    """Claude matntidan diagramma ma'lumotlarini ajratish"""
    # Format: [ДИАГРАММА: nomi | label1:val1, label2:val2, ...]
    pattern = r'\[(?:ДИАГРАММА|DIAGRAMMA|CHART|INFOGRAFIKA):\s*([^\|]+)\|([^\]]+)\]'
    matches = re.findall(pattern, claude_text, re.IGNORECASE)
    
    results = []
    for title, data_str in matches:
        entries = []
        for item in data_str.split(','):
            item = item.strip()
            if ':' in item:
                parts = item.split(':')
                label = parts[0].strip()
                try:
                    val = float(parts[1].strip().replace('%','').replace(',','.'))
                    entries.append((label, val))
                except: pass
        if entries:
            results.append({"title": title.strip(), "data": entries})
    return results

# ============================================================
# DIAGRAMMA CHIZISH FUNKSIYALARI
# ============================================================

def add_bar_chart(sl, data, title, x, y, w, h, accent_color):
    """Ustunli diagramma qo'shish"""
    try:
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE
        
        chart_data = ChartData()
        chart_data.categories = [d[0] for d in data[:8]]
        chart_data.add_series(title, [d[1] for d in data[:8]])
        
        chart = sl.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(x), Inches(y), Inches(w), Inches(h),
            chart_data
        ).chart
        
        # Rang berish
        try:
            plot = chart.plots[0]
            series = plot.series[0]
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = accent_color
        except: pass
        
        # Sarlavha
        try:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
            chart.chart_title.text_frame.paragraphs[0].font.bold = True
        except: pass
        
        # Legend o'chirish
        try: chart.has_legend = False
        except: pass
        
        return True
    except Exception as e:
        return False

def add_pie_chart(sl, data, title, x, y, w, h, accent_color):
    """Doiraviy diagramma qo'shish"""
    try:
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE
        
        chart_data = ChartData()
        chart_data.categories = [d[0] for d in data[:6]]
        chart_data.add_series(title, [d[1] for d in data[:6]])
        
        chart = sl.shapes.add_chart(
            XL_CHART_TYPE.PIE,
            Inches(x), Inches(y), Inches(w), Inches(h),
            chart_data
        ).chart
        
        try:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
        except: pass
        
        try:
            chart.has_legend = True
            chart.legend.position = 2  # bottom
            chart.legend.include_in_layout = False
        except: pass
        
        return True
    except: return False

def add_line_chart(sl, data, title, x, y, w, h, accent_color):
    """Chiziqli diagramma (trend)"""
    try:
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE
        
        chart_data = ChartData()
        chart_data.categories = [d[0] for d in data[:10]]
        chart_data.add_series(title, [d[1] for d in data[:10]])
        
        chart = sl.shapes.add_chart(
            XL_CHART_TYPE.LINE,
            Inches(x), Inches(y), Inches(w), Inches(h),
            chart_data
        ).chart
        
        try:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
        except: pass
        
        try:
            series = chart.plots[0].series[0]
            series.format.line.color.rgb = accent_color
            series.format.line.width = Pt(2.5)
        except: pass
        
        return True
    except: return False

def add_simple_bar_infographic(sl, data, title, x, y, w, h, accent_color, text_color):
    """Oddiy chiziqli infografika (chart kutubxonasisiz)"""
    try:
        max_val = max(d[1] for d in data) if data else 1
        bar_h = 0.3
        spacing = 0.45
        title_tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.35))
        tf = title_tb.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = text_color
        
        for i, (label, val) in enumerate(data[:6]):
            bar_y = y + 0.4 + i * spacing
            bar_w = (val / max_val) * (w - 2.5) if max_val > 0 else 0.1
            
            # Label
            label_tb = sl.shapes.add_textbox(Inches(x), Inches(bar_y), Inches(2.2), Inches(bar_h))
            lp = label_tb.text_frame.paragraphs[0]
            lp.text = label[:20]; lp.font.size = Pt(10); lp.font.color.rgb = text_color
            
            # Bar
            if bar_w > 0.05:
                bar = sl.shapes.add_shape(1, Inches(x+2.3), Inches(bar_y+0.03), Inches(bar_w), Inches(bar_h-0.08))
                bar.fill.solid(); bar.fill.fore_color.rgb = accent_color; bar.line.fill.background()
            
            # Value
            val_tb = sl.shapes.add_textbox(Inches(x+2.4+bar_w), Inches(bar_y), Inches(0.8), Inches(bar_h))
            vp = val_tb.text_frame.paragraphs[0]
            vp.text = f"{val:.1f}"; vp.font.size = Pt(10); vp.font.color.rgb = text_color
        
        return True
    except: return False

# ============================================================
# ASOSIY DIAGRAMMA QO'SHISH FUNKSIYASI
# ============================================================

def add_diagram_to_slide(sl, topic, slide_title, slide_text, diagram_data, 
                          accent_color, text_color, has_image=False):
    """Slaydga diagramma qo'shish"""
    if not diagram_data:
        return False
    
    d = diagram_data[0]
    title = d["title"]
    data = d["data"]
    
    if not data:
        return False
    
    # Joy hisoblash
    if has_image:
        x, y, w, h = 0.4, 4.2, 7.5, 2.8
    else:
        x, y, w, h = 0.4, 4.0, 12.5, 3.0
    
    dtype = detect_diagram_type(topic, slide_title, slide_text)
    
    # Avval python-pptx chart bilan sinab ko'ramiz
    success = False
    if dtype == "pie_chart" and len(data) >= 2:
        success = add_pie_chart(sl, data, title, x, y, w, h, accent_color)
    elif dtype == "line_chart" and len(data) >= 2:
        success = add_line_chart(sl, data, title, x, y, w, h, accent_color)
    
    if not success:
        # Fallback: oddiy chiziqli infografika
        success = add_simple_bar_infographic(sl, data, title, x, y, w, h, accent_color, text_color)
    
    return success



# ============================================================
# CLAUDE API
# ============================================================
def claude(prompt, system="", max_tok=4000, model=None):
    if not CLAUDE_API_KEY:
        return "Claude API sozlanmagan!"
    if model is None:
        model = HAIKU_MODEL
    try:
        r = requests.post(
            "https://api.anthropic.com/v1/messages",
            headers={
                "x-api-key": CLAUDE_API_KEY,
                "anthropic-version": "2023-06-01",
                "content-type": "application/json"
            },
            json={
                "model": model,
                "max_tokens": max_tok,
                "system": system,
                "messages": [{"role": "user", "content": prompt}]
            },
            timeout=180
        )
        if r.status_code == 200:
            return r.json()["content"][0]["text"]
        else:
            logger.error(f"Claude API error: {r.status_code} {r.text}")
            return f"API xatosi: {r.status_code}"
    except Exception as e:
        logger.error(f"Claude xato: {e}")
        return f"Xatolik: {e}"

def clean_text(text):
    """AI javobidan keraksiz belgilarni tozalash"""
    import re as _re
    text = _re.sub(r'\*\*(.+?)\*\*', r'\1', text, flags=_re.DOTALL)
    text = _re.sub(r'\*(.+?)\*', r'\1', text)
    text = _re.sub(r'#{1,6}\s*', '', text)
    text = _re.sub(r'[`~]', '', text)
    text = _re.sub(r'\n{3,}', '\n\n', text)
    lines = []
    for line in text.split('\n'):
        line = line.strip()
        if not (line.upper().startswith('SLAYD') or
                line.upper().startswith('SLIDE') or
                line.upper().startswith('INFOGRAFIKA')):
            line = _re.sub(r'^[-•►▸*]+\s*', '', line)
        lines.append(line)
    return '\n'.join(lines).strip()

# ============================================================
# MAZMUN YARATISH
# ============================================================
def gen_prez(topic, slides, lang, ud={}, plans=5):
    """Prezentatsiya mazmuni yaratish"""
    ln = LN.get(lang, "o'zbek")
    info = build_info(ud)
    subject = f"\nFan: {ud['subject']}" if ud.get('subject') else ""
    book = f"\nManba: {ud['book_name']}" if ud.get('book_name') else ""
    plans_count = ud.get("plans_count", plans)
    with_diagram = ud.get("with_diagram", False)

    # Web search
    web_ctx = ""
    try:
        web_info = web_search_topic(f"{topic} ma'lumotlar faktlar", lang)
        if web_info: web_ctx = f"\n\nQo'shimcha ma'lumot:\n{web_info[:800]}"
    except: pass

    # Bo'limlar ro'yxatini yaratish
    bolimlar = []
    for i in range(1, plans_count + 1):
        bolimlar.append(f"{i}. [Bo'lim {i} sarlavhasi]")
    bolimlar_str = "\n".join(bolimlar)

    # Slaydlar ro'yxati
    slide_plan = f"SLAYD 1: {topic}\n"
    slide_plan += f"SLAYD 2: REJA\n"
    slayd_num = 3
    bolim_per_slide = max(1, (slides - 3) // plans_count)
    for b in range(1, plans_count + 1):
        slide_plan += f"SLAYD {slayd_num}: [Bo'lim {b} nomi]\n"
        slayd_num += 1
        if slayd_num >= slides:
            break
    slide_plan += f"SLAYD {slides}: Xulosa\n"

    diag_instruction = ""
    if with_diagram:
        diag_instruction = (
            f"\nDIAGRAMMA: Agar mavzuda statistik ma'lumotlar bo'lsa, "
            f"faqat shu slaydda [DIAGRAMMA: nom | label1:son1, label2:son2] qo'sh. "
            f"Diagramma MATN PASTIDA joylashtiriladi."
        )

    prompt = f"""Mavzu: {topic}
Slaydlar soni: {slides} ta
Bo'limlar soni: {plans_count} ta
Til: {ln}
{info}{subject}{book}{web_ctx}

TOPSHIRIQ: Quyidagi {slides} ta slaydni to'liq yoz:

{slide_plan}

MUHIM QOIDALAR:
1. Har slayd AYNAN "SLAYD N:" bilan boshlansin
2. Jami {slides} ta slayd yoz — kam yozma!
3. SLAYD 2 (REJA) da:
   - "REJA" so'zi katta harfda MARKAZDA
   - Quyida tartibli ro'yxat:
{bolimlar_str}
   - Har bo'lim haqida 1-2 jumlali qisqa tavsif
4. SLAYD 3 dan {slides-1} gacha — Reja bo'limlari tartibida, har birida 250-300 so'z
5. **, ##, # belgisi TAQIQLANGAN{diag_instruction}"""

    system = f"""Sen {ln} tilida professional prezentatsiya yozuvchisan.
QOIDA 1: Jami {slides} ta slayd yoz — bu MAJBURIY, kamroq yozma!
QOIDA 2: Har slayd "SLAYD N:" bilan boshlansin.
QOIDA 3: SLAYD 2 da REJA so'zi va tartibli bo'limlar bo'lsin.
QOIDA 4: Har slaydda 250-300 so'z bo'lsin.
QOIDA 5: Markdown belgisi ishlatma."""

    # Bitta so'rovda yaratish
    result = claude(prompt, system, 8000, model=SONNET_MODEL)

    if not result or "API xatosi" in result or "timeout" in result.lower():
        logger.error(f"Claude failed in gen_prez: {result[:100]}")
        fallback = ""
        for i in range(1, slides + 1):
            fallback += f"SLAYD {i}: {topic if i == 1 else ('Reja' if i == 2 else f'Bo\'lim {i-2}')}\n"
            fallback += f"{topic} haqida ma'lumot.\n\n"
        return fallback

    result = result.replace("**", "").replace("## ", "").replace("# ", "").replace("##", "")
    count = len(re.findall(r'SLAYD\s*\d+\s*:', result, re.IGNORECASE))
    logger.info(f"PREZ: {count} slayd topildi (so'ralgan: {slides})")
    return result


def gen_doc(svc, topic, pages, lang, ud={}):
    """Hujjat yaratish (referat, kurs ishi, maqola, mustaqil)"""
    ln = LN.get(lang, "o'zbek")
    info = build_info(ud)
    wpg = {"referat": 300, "kurs": 350, "mustaqil": 280, "maqola": 320}.get(svc, 300)
    subject = f"\nFan: {ud['subject']}" if ud.get('subject') else ""

    structs = {
        "referat": (
            "KIRISH (200+ so'z)\n"
            "I BOB — Nazariy asoslar (300+ so'z)\n"
            "II BOB — Tahlil va muhokama (300+ so'z)\n"
            "III BOB — Amaliy qo'llanish (300+ so'z)\n"
            "XULOSA (150+ so'z)\n"
            "FOYDALANILGAN ADABIYOTLAR (10 ta manba)"
        ),
        "kurs": (
            "MUNDARIJA\n"
            "KIRISH — maqsad, vazifalar, dolzarblik (300+ so'z)\n"
            "I BOB — Nazariy asos (400+ so'z)\n"
            "II BOB — Tahlil va tadqiqot (400+ so'z)\n"
            "III BOB — Tavsiyalar va xulosalar (300+ so'z)\n"
            "XULOSA (200+ so'z)\n"
            "FOYDALANILGAN ADABIYOTLAR (15 ta APA)\n"
            "ILOVALAR"
        ),
        "mustaqil": (
            "KIRISH (200+ so'z)\n"
            "ASOSIY QISM 1 (300+ so'z)\n"
            "ASOSIY QISM 2 (300+ so'z)\n"
            "XULOSA (150+ so'z)\n"
            "FOYDALANILGAN ADABIYOTLAR (8 ta)"
        ),
        "maqola": (
            "ANNOTATSIYA (150 so'z)\n"
            "KALIT SO'ZLAR (7 ta)\n"
            "ABSTRACT (inglizcha, 150 so'z)\n"
            "KEYWORDS (7 ta)\n"
            "KIRISH (300+ so'z)\n"
            "ADABIYOTLAR TAHLILI (250+ so'z)\n"
            "METODOLOGIYA (200+ so'z)\n"
            "NATIJALAR VA MUHOKAMA (400+ so'z)\n"
            "XULOSA (200+ so'z)\n"
            "FOYDALANILGAN ADABIYOTLAR (15 ta APA)"
        )
    }
    names = {"referat": "referat", "kurs": "kurs ishi", "mustaqil": "mustaqil ish", "maqola": "ilmiy maqola"}
    struct = structs.get(svc, structs["referat"])
    use_model = SONNET_MODEL if svc in ("kurs", "maqola") else HAIKU_MODEL

    prompt = (
        f"Mavzu: {topic}\n"
        f"Hajm: {pages} bet ({pages * wpg}+ so'z)\n"
        f"{info}{subject}\n\n"
        f"{ln} tilida TO'LIQ {names.get(svc, 'hujjat')} yozing:\n\n"
        f"{struct}\n\n"
        f"QATIY TALABLAR:\n"
        f"1. Faqat ilmiy kitoblar va tasdiqlangan manbalardan\n"
        f"2. Har bo'limda aniq faktlar, raqamlar, foizlar\n"
        f"3. Hech qanday **, *, #, ` belgisi ishlatilmasin\n"
        f"4. Imlo va grammatika 100% to'g'ri\n"
        f"5. Har bo'lim to'liq, qisqartirma yo'q\n"
        f"6. Oxirida foydalanilgan adabiyotlar ro'yxati"
    )
    system = (
        f"Sen O'zbekistonning eng professional {ln} akademik yozuvchisisan. "
        f"Faqat ilmiy manbalardan foydalanasan. "
        f"Markdown belgisi ASLO ishlatmaysan. "
        f"Imlo va grammatika xatosiz yozasan."
    )
    # Manba qo'shish
    source_ctx = build_prompt_with_source(svc, topic, pages, lang, ud)
    if source_ctx:
        prompt = prompt + source_ctx
    
    result = claude(prompt, system, min(pages * 200 + 2000, 4000), model=use_model)
    return clean_text(result)

def gen_test(topic, count, lang):
    """Test savollari yaratish"""
    ln = LN.get(lang, "o'zbek")
    prompt = (
        f"Mavzu: {topic}\nSavollar soni: {count}\nTil: {ln}\n\n"
        f"{count} ta professional test savoli yarating.\n\n"
        f"FORMAT (qat'iy shu formatda):\n"
        f"1. [Savol matni]\n"
        f"A) [javob]\nB) [javob]\nC) [javob]\nD) [javob]\n"
        f"To'g'ri javob: A\n\n"
        f"2. [Savol matni]\n"
        f"...\n\n"
        f"TALABLAR:\n"
        f"1. Hech qanday **, *, # belgisi yo'q\n"
        f"2. Imlo xatosiz\n"
        f"3. Savollar turli murakkablik darajasida\n"
        f"4. To'g'ri javob aniq ko'rsatilsin"
    )
    system = f"Sen professional {ln} test yaratuvchisisan. Markdown ishlatmaysan. Imlo xatosiz."
    result = claude(prompt, system, min(count * 100, 4000))
    return clean_text(result)

def fix_spell(text, lang):
    """Imlo tuzatish"""
    ln = LN.get(lang, "o'zbek")
    result = claude(
        f"{ln} tilida imlo tuzat. Faqat tuzatilgan matnni qaytар:\n\n{text[:3000]}",
        f"Sen {ln} imlo mutaxassisisanm. Faqat tuzatilgan matn.",
        3500
    )
    return clean_text(result)

# ============================================================
# RASM OLISH
# ============================================================

def web_search_topic(query, lang="uz"):
    """Mavzu bo'yicha internetdan ma'lumot qidirish"""
    try:
        headers = {"User-Agent": "Mozilla/5.0"}
        # DuckDuckGo API
        r = requests.get(
            f"https://api.duckduckgo.com/?q={requests.utils.quote(query)}&format=json&no_html=1&skip_disambig=1",
            headers=headers, timeout=10
        )
        if r.status_code == 200:
            data = r.json()
            abstract = data.get("AbstractText", "")
            related = [t.get("Text","") for t in data.get("RelatedTopics", [])[:5] if t.get("Text")]
            if abstract or related:
                info = abstract + "\n" + "\n".join(related)
                return info[:2000]
    except Exception as e:
        logger.warning(f"Web search: {e}")
    return ""

def build_prompt_with_source(svc, topic, pages, lang, ud):
    """Manba asosida prompt yaratish"""
    source_type = ud.get("source_type", "none")
    source_text = ud.get("source_text", "")
    book_name = ud.get("book_name", "")
    
    # Kitob nomi yoki matn berilgan bo'lsa
    if source_type == "text" and book_name:
        web_info = web_search_topic(f"{book_name} {topic}", lang)
        source_context = f"\n\nManba kitob: {book_name}\n"
        if web_info:
            source_context += f"Kitobdan ma'lumotlar:\n{web_info[:1500]}\n"
        source_context += f"\nFAQAT shu kitob/{book_name} ma'lumotlaridan foydalaning!"
    elif source_type == "pdf" and source_text:
        source_context = f"\n\nQuyidagi kitob matni asosida yarating:\n{source_text[:3000]}\n\nFaqat berilgan matn asosida yozing!"
    else:
        web_info = web_search_topic(f"{topic} statistics data", lang)
        source_context = f"\n\nQo'shimcha ma'lumot:\n{web_info[:1000]}" if web_info else ""
    
    return source_context

def get_image(query):
    """Internetdan mavzuga mos rasm olish — Pixabay API"""
    try:
        import hashlib, time
        # Mavzuni inglizchaga tarjima qilish
        try:
            en_q = claude(f"Translate to 2-3 English words only (no explanation): {query}", max_tok=15)
            en_q = en_q.strip().split("\n")[0][:40]
        except:
            en_q = query[:40]
        
        search_q = requests.utils.quote(en_q)
        
        # 1. Pixabay (bepul, API key shart emas)
        pixabay_key = os.environ.get("PIXABAY_KEY", "")
        if pixabay_key:
            r = requests.get(
                f"https://pixabay.com/api/?key={pixabay_key}&q={search_q}"
                f"&image_type=photo&orientation=horizontal&min_width=800&per_page=5&safesearch=true",
                timeout=10)
            if r.status_code == 200:
                hits = r.json().get("hits", [])
                if hits:
                    import random
                    hit = random.choice(hits[:3])
                    img_url = hit.get("webformatURL", "")
                    if img_url:
                        ir = requests.get(img_url, timeout=15)
                        if ir.status_code == 200 and len(ir.content) > 5000:
                            buf = BytesIO(ir.content); buf.seek(0)
                            return buf
        
        # 2. Unsplash (agar kalit bo'lsa)
        if UNSPLASH_KEY:
            headers = {"Authorization": f"Client-ID {UNSPLASH_KEY}"}
            r = requests.get(
                f"https://api.unsplash.com/search/photos?query={search_q}&per_page=5&orientation=landscape",
                headers=headers, timeout=10)
            if r.status_code == 200:
                results = r.json().get("results", [])
                if results:
                    img_url = results[0].get("urls", {}).get("regular", "")
                    if img_url:
                        ir = requests.get(img_url, timeout=15)
                        if ir.status_code == 200 and len(ir.content) > 5000:
                            buf = BytesIO(ir.content); buf.seek(0)
                            return buf
        
        # 3. Fallback - rasm topilmasa None qaytarish
        return None
        
    except Exception as e:
        logger.warning(f"Rasm olishda xato: {e}")
    return None

# ============================================================
# 42 SHABLON
# ============================================================

# ============================================================
# 30 TA PROFESSIONAL SHABLON
# ============================================================
# ============================================================
# 30 TA PROFESSIONAL SHABLON MODULI
# Har biri to'liq boshqacha dizayn
# ============================================================
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import math

def rgb(r, g, b):
    return RGBColor(r, g, b)

# ============================================================
# 30 TA SHABLON TA'RIFI
# Har birida: name, category, bg_func, decor_func, title_style, text_style
# ============================================================

TEMPLATES_30 = {
    # ── MINIMALIST GURUH (1-6) ─────────────────────────────
    "1": {
        "name": "⬜ Toza Oq",
        "category": "minimalist",
        "preview": "🟫",
        "desc": "Toza oq fon, ingichka chiziqlar",
        "bg1": (255,255,255), "bg2": (248,249,250),
        "accent": (33,150,243), "title_c": (15,15,15),
        "text_c": (40,40,40), "line_c": (33,150,243),
        "title_size": 32, "text_size": 15,
        "style": "minimalist_clean"
    },
    "2": {
        "name": "🔵 Ko'k Chiziq",
        "category": "minimalist",
        "preview": "🟦",
        "desc": "Oq fon, ko'k aksent chiziqlari",
        "bg1": (252,253,255), "bg2": (240,247,255),
        "accent": (25,118,210), "title_c": (13,71,161),
        "text_c": (30,30,50), "line_c": (25,118,210),
        "title_size": 30, "text_size": 14,
        "style": "minimalist_line"
    },
    "3": {
        "name": "⚫ Qora Elegant",
        "category": "minimalist",
        "preview": "⬛",
        "desc": "Qora fon, oltin detallar",
        "bg1": (18,18,18), "bg2": (28,28,28),
        "accent": (212,175,55), "title_c": (255,255,255),
        "text_c": (220,220,220), "line_c": (212,175,55),
        "title_size": 32, "text_size": 15,
        "style": "dark_elegant"
    },
    "4": {
        "name": "🌫️ Kulrang Sof",
        "category": "minimalist",
        "preview": "🔘",
        "desc": "Kulrang gradient, zamonaviy ko'rinish",
        "bg1": (236,239,241), "bg2": (207,216,220),
        "accent": (84,110,122), "title_c": (33,33,33),
        "text_c": (55,55,55), "line_c": (84,110,122),
        "title_size": 30, "text_size": 14,
        "style": "minimalist_gray"
    },
    "5": {
        "name": "🟢 Yashil Toza",
        "category": "minimalist",
        "preview": "🟩",
        "desc": "Yashil aksent, toza minimalist",
        "bg1": (255,255,255), "bg2": (232,245,233),
        "accent": (56,142,60), "title_c": (27,94,32),
        "text_c": (33,33,33), "line_c": (56,142,60),
        "title_size": 30, "text_size": 14,
        "style": "minimalist_green"
    },
    "6": {
        "name": "🔴 Qizil Bold",
        "category": "minimalist",
        "preview": "🟥",
        "desc": "Qizil aksent, kuchli zamonaviy",
        "bg1": (255,255,255), "bg2": (255,235,238),
        "accent": (211,47,47), "title_c": (183,28,28),
        "text_c": (33,33,33), "line_c": (211,47,47),
        "title_size": 30, "text_size": 14,
        "style": "minimalist_red"
    },

    # ── GRADIENT GURUH (7-12) ─────────────────────────────
    "7": {
        "name": "🌊 Okean Gradient",
        "category": "gradient",
        "preview": "🌊",
        "desc": "Ko'k-moviy gradient",
        "bg1": (13,71,161), "bg2": (3,169,244),
        "accent": (255,255,255), "title_c": (255,255,255),
        "text_c": (224,247,250), "line_c": (255,255,255),
        "title_size": 32, "text_size": 15,
        "style": "gradient_ocean"
    },
    "8": {
        "name": "🌅 Quyosh Botishi",
        "category": "gradient",
        "preview": "🌅",
        "desc": "To'q sariq-qizil gradient",
        "bg1": (183,28,28), "bg2": (255,143,0),
        "accent": (255,255,255), "title_c": (255,255,255),
        "text_c": (255,243,224), "line_c": (255,255,255),
        "title_size": 32, "text_size": 15,
        "style": "gradient_sunset"
    },
    "9": {
        "name": "💜 Binafsha Tun",
        "category": "gradient",
        "preview": "💜",
        "desc": "To'q binafsha gradient",
        "bg1": (49,27,146), "bg2": (123,31,162),
        "accent": (206,147,216), "title_c": (255,255,255),
        "text_c": (237,231,246), "line_c": (206,147,216),
        "title_size": 32, "text_size": 15,
        "style": "gradient_purple"
    },
    "10": {
        "name": "🌿 O'rmon Gradient",
        "category": "gradient",
        "preview": "🌿",
        "desc": "Yashil gradient tabiat uslubi",
        "bg1": (27,94,32), "bg2": (100,181,46),
        "accent": (255,255,255), "title_c": (255,255,255),
        "text_c": (232,245,233), "line_c": (255,255,255),
        "title_size": 32, "text_size": 15,
        "style": "gradient_forest"
    },
    "11": {
        "name": "🌌 Kosmik",
        "category": "gradient",
        "preview": "🌌",
        "desc": "Qora-ko'k kosmik gradient",
        "bg1": (5,5,30), "bg2": (30,30,80),
        "accent": (100,181,246), "title_c": (255,255,255),
        "text_c": (200,220,255), "line_c": (100,181,246),
        "title_size": 32, "text_size": 15,
        "style": "gradient_cosmic"
    },
    "12": {
        "name": "🍑 Shaftoli",
        "category": "gradient",
        "preview": "🍑",
        "desc": "Pushti-shaftoli yumshoq gradient",
        "bg1": (255,138,101), "bg2": (255,193,157),
        "accent": (255,255,255), "title_c": (255,255,255),
        "text_c": (255,243,240), "line_c": (255,255,255),
        "title_size": 30, "text_size": 14,
        "style": "gradient_peach"
    },

    # ── GEOMETRIC GURUH (13-18) ───────────────────────────
    "13": {
        "name": "🔷 Ko'k Geometrik",
        "category": "geometric",
        "preview": "🔷",
        "desc": "Ko'k uchburchaklar va shakllar",
        "bg1": (255,255,255), "bg2": (227,242,253),
        "accent": (21,101,192), "title_c": (13,71,161),
        "text_c": (25,25,25), "line_c": (21,101,192),
        "title_size": 30, "text_size": 14,
        "style": "geometric_blue"
    },
    "14": {
        "name": "🔶 Oltin Geometrik",
        "category": "geometric",
        "preview": "🔶",
        "desc": "Oltin-qora geometrik uslub",
        "bg1": (25,25,25), "bg2": (40,35,10),
        "accent": (255,196,0), "title_c": (255,196,0),
        "text_c": (240,240,240), "line_c": (255,196,0),
        "title_size": 30, "text_size": 14,
        "style": "geometric_gold"
    },
    "15": {
        "name": "🔺 Qizil Uchburchak",
        "category": "geometric",
        "preview": "🔺",
        "desc": "Qizil geometrik dizayn",
        "bg1": (255,255,255), "bg2": (255,235,238),
        "accent": (198,40,40), "title_c": (183,28,28),
        "text_c": (30,30,30), "line_c": (198,40,40),
        "title_size": 30, "text_size": 14,
        "style": "geometric_red"
    },
    "16": {
        "name": "💠 Moviy Mozaika",
        "category": "geometric",
        "preview": "💠",
        "desc": "Moviy mozaika geometrik",
        "bg1": (2,136,209), "bg2": (1,87,155),
        "accent": (255,255,255), "title_c": (255,255,255),
        "text_c": (225,245,254), "line_c": (255,255,255),
        "title_size": 30, "text_size": 14,
        "style": "geometric_mosaic"
    },
    "17": {
        "name": "🟫 Bronza Geometrik",
        "category": "geometric",
        "preview": "🟫",
        "desc": "Bronza-jigarrang geometrik",
        "bg1": (62,39,35), "bg2": (78,52,46),
        "accent": (188,143,143), "title_c": (255,204,128),
        "text_c": (255,224,178), "line_c": (188,143,143),
        "title_size": 30, "text_size": 14,
        "style": "geometric_bronze"
    },
    "18": {
        "name": "⬡ Oltiburchak",
        "category": "geometric",
        "preview": "⬡",
        "desc": "Zamonaviy oltiburchak pattern",
        "bg1": (245,245,245), "bg2": (224,224,224),
        "accent": (97,97,97), "title_c": (33,33,33),
        "text_c": (50,50,50), "line_c": (97,97,97),
        "title_size": 28, "text_size": 13,
        "style": "geometric_hex"
    },

    # ── KORPORATIV GURUH (19-24) ──────────────────────────
    "19": {
        "name": "🏢 Navy Professional",
        "category": "corporate",
        "preview": "🏢",
        "desc": "Navy ko'k korporativ uslub",
        "bg1": (255,255,255), "bg2": (232,234,246),
        "accent": (26,35,126), "title_c": (26,35,126),
        "text_c": (33,33,33), "line_c": (26,35,126),
        "title_size": 30, "text_size": 14,
        "style": "corporate_navy"
    },
    "20": {
        "name": "🌐 Temir Korporativ",
        "category": "corporate",
        "preview": "🌐",
        "desc": "Temir-kulrang professional",
        "bg1": (33,33,33), "bg2": (55,55,55),
        "accent": (96,125,139), "title_c": (255,255,255),
        "text_c": (200,200,200), "line_c": (96,125,139),
        "title_size": 30, "text_size": 14,
        "style": "corporate_steel"
    },
    "21": {
        "name": "📊 Ma'lumot Tahlil",
        "category": "corporate",
        "preview": "📊",
        "desc": "Ma'lumot va tahlil uchun maxsus",
        "bg1": (250,251,252), "bg2": (236,240,241),
        "accent": (41,128,185), "title_c": (44,62,80),
        "text_c": (52,73,94), "line_c": (41,128,185),
        "title_size": 28, "text_size": 13,
        "style": "corporate_data"
    },
    "22": {
        "name": "🎯 Maqsad",
        "category": "corporate",
        "preview": "🎯",
        "desc": "Sariq-qora biznes uslub",
        "bg1": (255,255,255), "bg2": (255,253,231),
        "accent": (245,127,23), "title_c": (230,81,0),
        "text_c": (33,33,33), "line_c": (245,127,23),
        "title_size": 30, "text_size": 14,
        "style": "corporate_target"
    },
    "23": {
        "name": "💼 Klassik Biznes",
        "category": "corporate",
        "preview": "💼",
        "desc": "Klassik oq-ko'k biznes",
        "bg1": (255,255,255), "bg2": (245,247,250),
        "accent": (0,90,160), "title_c": (0,70,127),
        "text_c": (40,40,60), "line_c": (0,90,160),
        "title_size": 32, "text_size": 15,
        "style": "corporate_classic"
    },
    "24": {
        "name": "🔬 Ilmiy Tadqiqot",
        "category": "corporate",
        "preview": "🔬",
        "desc": "Ilmiy-akademik professional uslub",
        "bg1": (252,252,252), "bg2": (240,242,245),
        "accent": (0,121,107), "title_c": (0,96,100),
        "text_c": (30,40,40), "line_c": (0,121,107),
        "title_size": 28, "text_size": 13,
        "style": "corporate_academic"
    },

    # ── IJODIY GURUH (25-30) ──────────────────────────────
    "25": {
        "name": "🎨 Ijodiy Rang",
        "category": "creative",
        "preview": "🎨",
        "desc": "Rangli ijodiy dizayn",
        "bg1": (255,255,255), "bg2": (250,250,255),
        "accent": (103,58,183), "title_c": (81,45,168),
        "text_c": (30,30,50), "line_c": (103,58,183),
        "title_size": 30, "text_size": 14,
        "style": "creative_colorful"
    },
    "26": {
        "name": "🌸 Pushti Zamonaviy",
        "category": "creative",
        "preview": "🌸",
        "desc": "Pushti-oq zamonaviy",
        "bg1": (255,255,255), "bg2": (252,228,236),
        "accent": (216,27,96), "title_c": (173,20,87),
        "text_c": (33,33,33), "line_c": (216,27,96),
        "title_size": 30, "text_size": 14,
        "style": "creative_pink"
    },
    "27": {
        "name": "🔥 Olov",
        "category": "creative",
        "preview": "🔥",
        "desc": "Qizil-sariq olov effekti",
        "bg1": (15,15,15), "bg2": (40,10,0),
        "accent": (255,87,34), "title_c": (255,193,7),
        "text_c": (255,224,178), "line_c": (255,87,34),
        "title_size": 32, "text_size": 15,
        "style": "creative_fire"
    },
    "28": {
        "name": "🌊 To'lqin",
        "category": "creative",
        "preview": "🌊",
        "desc": "To'lqin effektli zamonaviy",
        "bg1": (0,150,136), "bg2": (0,105,92),
        "accent": (255,255,255), "title_c": (255,255,255),
        "text_c": (224,242,241), "line_c": (255,255,255),
        "title_size": 30, "text_size": 14,
        "style": "creative_wave"
    },
    "29": {
        "name": "✨ Yulduz Kechasi",
        "category": "creative",
        "preview": "✨",
        "desc": "Yulduzli tun osmoni",
        "bg1": (10,10,40), "bg2": (30,30,70),
        "accent": (255,215,0), "title_c": (255,255,200),
        "text_c": (200,210,255), "line_c": (255,215,0),
        "title_size": 30, "text_size": 14,
        "style": "creative_starnight"
    },
    "30": {
        "name": "🎓 Akademik Klassik",
        "category": "creative",
        "preview": "🎓",
        "desc": "Akademik klassik uslub",
        "bg1": (250,245,230), "bg2": (240,230,210),
        "accent": (120,90,30), "title_c": (80,50,10),
        "text_c": (60,40,20), "line_c": (120,90,30),
        "title_size": 30, "text_size": 14,
        "style": "academic_classic"
    },
}

# ============================================================
# SHABLON FONGA CHIZISH FUNKSIYALARI
# ============================================================

def draw_minimalist_clean(sl, tmpl):
    """Toza minimalist - ingichka chiziqlar"""
    acc = rgb(*tmpl.get("accent", (33,150,243)))
    # Pastki chiziq
    try:
        bar = sl.shapes.add_shape(1, Inches(0), Inches(6.9), Inches(13.33), Inches(0.1))
        bar.fill.solid(); bar.fill.fore_color.rgb = acc; bar.line.fill.background()
    except: pass
    # Yuqori ingichka chiziq
    try:
        bar2 = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.05))
        bar2.fill.solid(); bar2.fill.fore_color.rgb = acc; bar2.line.fill.background()
    except: pass

def draw_minimalist_line(sl, tmpl):
    """Ko'k chiziq - chapda vertikal chiziq"""
    acc = rgb(*tmpl.get("accent", (33,150,243)))
    try:
        bar = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.12), Inches(7.5))
        bar.fill.solid(); bar.fill.fore_color.rgb = acc; bar.line.fill.background()
    except: pass
    try:
        bar2 = sl.shapes.add_shape(1, Inches(0), Inches(6.9), Inches(13.33), Inches(0.08))
        bar2.fill.solid(); bar2.fill.fore_color.rgb = acc; bar2.line.fill.background()
    except: pass

def draw_dark_elegant(sl, tmpl):
    """Qora elegant - oltin chiziqlar"""
    acc = rgb(*tmpl.get("accent", (33,150,243)))
    try:
        # Yuqori chiziq
        bar = sl.shapes.add_shape(1, Inches(0.3), Inches(0.05), Inches(12.7), Inches(0.06))
        bar.fill.solid(); bar.fill.fore_color.rgb = acc; bar.line.fill.background()
        # Pastki chiziq
        bar2 = sl.shapes.add_shape(1, Inches(0.3), Inches(7.35), Inches(12.7), Inches(0.06))
        bar2.fill.solid(); bar2.fill.fore_color.rgb = acc; bar2.line.fill.background()
        # O'ng doira
        c = sl.shapes.add_shape(9, Inches(11.5), Inches(5.5), Inches(2.5), Inches(2.5))
        c.fill.solid(); c.fill.fore_color.rgb = acc
        c.fill.fore_color.transparency = 0.9; c.line.fill.background()
    except: pass

def draw_gradient_ocean(sl, tmpl):
    """Okean gradient - to'lqin shakli"""
    try:
        # O'ng pastda katta doira
        c1 = sl.shapes.add_shape(9, Inches(9), Inches(4), Inches(6), Inches(6))
        c1.fill.solid(); c1.fill.fore_color.rgb = rgb(255,255,255)
        c1.fill.fore_color.transparency = 0.9; c1.line.fill.background()
        # Chap tomonda kichik doiralar
        c2 = sl.shapes.add_shape(9, Inches(-1), Inches(-0.5), Inches(3), Inches(3))
        c2.fill.solid(); c2.fill.fore_color.rgb = rgb(255,255,255)
        c2.fill.fore_color.transparency = 0.85; c2.line.fill.background()
    except: pass

def draw_geometric_blue(sl, tmpl):
    """Ko'k geometrik - uchburchaklar"""
    acc = rgb(*tmpl.get("accent", (33,150,243)))
    try:
        # O'ng uchda uchburchak
        from pptx.util import Pt as _Pt
        tri = sl.shapes.add_shape(5, Inches(10), Inches(0), Inches(3.33), Inches(3))
        tri.fill.solid(); tri.fill.fore_color.rgb = acc
        tri.fill.fore_color.transparency = 0.85; tri.line.fill.background()
        # Pastda kichik uchburchak
        tri2 = sl.shapes.add_shape(5, Inches(0), Inches(5.5), Inches(2), Inches(2))
        tri2.fill.solid(); tri2.fill.fore_color.rgb = acc
        tri2.fill.fore_color.transparency = 0.8; tri2.line.fill.background()
        # Chiziq
        bar = sl.shapes.add_shape(1, Inches(0), Inches(7.3), Inches(13.33), Inches(0.08))
        bar.fill.solid(); bar.fill.fore_color.rgb = acc; bar.line.fill.background()
    except: pass

def draw_geometric_gold(sl, tmpl):
    """Oltin geometrik - diagonal shakllar"""
    acc = rgb(*tmpl.get("accent", (33,150,243)))
    try:
        # Diagonal chiziqlar
        for i in range(3):
            bar = sl.shapes.add_shape(1, Inches(10+i*0.3), Inches(0), Inches(0.08), Inches(7.5))
            bar.fill.solid(); bar.fill.fore_color.rgb = acc
            bar.fill.fore_color.transparency = 0.7; bar.line.fill.background()
        # Pastki chiziq
        bar2 = sl.shapes.add_shape(1, Inches(0), Inches(7.2), Inches(13.33), Inches(0.06))
        bar2.fill.solid(); bar2.fill.fore_color.rgb = acc; bar2.line.fill.background()
    except: pass

def draw_geometric_mosaic(sl, tmpl):
    """Mozaika - kvadratlar"""
    acc = rgb(255,255,255)
    try:
        for i in range(4):
            for j in range(2):
                sq = sl.shapes.add_shape(1, Inches(10.5+i*0.7), Inches(j*0.7), Inches(0.6), Inches(0.6))
                sq.fill.solid(); sq.fill.fore_color.rgb = acc
                sq.fill.fore_color.transparency = 0.8+j*0.05; sq.line.fill.background()
    except: pass

def draw_corporate_navy(sl, tmpl):
    """Navy korporativ - chapda rangli panel"""
    acc = rgb(*tmpl.get("accent", (33,150,243)))
    try:
        # Chapda tor panel
        panel = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.8), Inches(7.5))
        panel.fill.solid(); panel.fill.fore_color.rgb = acc; panel.line.fill.background()
        # Sarlavha osti chizig'i
        bar = sl.shapes.add_shape(1, Inches(0.9), Inches(1.5), Inches(11), Inches(0.05))
        bar.fill.solid(); bar.fill.fore_color.rgb = acc; bar.line.fill.background()
    except: pass

def draw_corporate_steel(sl, tmpl):
    """Temir - metallik ko'rinish"""
    acc = rgb(*tmpl.get("accent", (33,150,243)))
    try:
        bar = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.15))
        bar.fill.solid(); bar.fill.fore_color.rgb = acc; bar.line.fill.background()
        bar2 = sl.shapes.add_shape(1, Inches(0), Inches(7.35), Inches(13.33), Inches(0.08))
        bar2.fill.solid(); bar2.fill.fore_color.rgb = acc; bar2.line.fill.background()
    except: pass

def draw_creative_fire(sl, tmpl):
    """Olov - qizil-sariq shakllar"""
    try:
        c1 = sl.shapes.add_shape(9, Inches(10), Inches(3), Inches(5), Inches(5))
        c1.fill.solid(); c1.fill.fore_color.rgb = rgb(255,87,34)
        c1.fill.fore_color.transparency = 0.85; c1.line.fill.background()
        c2 = sl.shapes.add_shape(9, Inches(-1), Inches(4), Inches(3), Inches(3))
        c2.fill.solid(); c2.fill.fore_color.rgb = rgb(255,193,7)
        c2.fill.fore_color.transparency = 0.8; c2.line.fill.background()
    except: pass

def draw_creative_wave(sl, tmpl):
    """To'lqin - yumaloq shakllar"""
    try:
        for i, (x, y, s) in enumerate([(11,5,4),(-0.5,5.5,3),(6,6.5,2)]):
            c = sl.shapes.add_shape(9, Inches(x), Inches(y), Inches(s), Inches(s))
            c.fill.solid(); c.fill.fore_color.rgb = rgb(255,255,255)
            c.fill.fore_color.transparency = 0.85; c.line.fill.background()
    except: pass

def draw_starnight(sl, tmpl):
    """Yulduzli tun - kichik doiralar"""
    acc = rgb(*tmpl.get("accent", (33,150,243)))
    import random
    random.seed(42)
    try:
        for _ in range(12):
            x = random.uniform(0, 12)
            y = random.uniform(0, 7)
            s = random.uniform(0.05, 0.12)
            c = sl.shapes.add_shape(9, Inches(x), Inches(y), Inches(s), Inches(s))
            c.fill.solid(); c.fill.fore_color.rgb = acc
            c.fill.fore_color.transparency = random.uniform(0.3, 0.7)
            c.line.fill.background()
    except: pass

def draw_academic(sl, tmpl):
    """Akademik - ramka"""
    acc = rgb(*tmpl.get("accent", (33,150,243)))
    try:
        for coords in [(Inches(0.2), Inches(0.2), Inches(12.9), Inches(0.05)),
                       (Inches(0.2), Inches(7.2), Inches(12.9), Inches(0.05)),
                       (Inches(0.2), Inches(0.2), Inches(0.05), Inches(7.1)),
                       (Inches(13.1), Inches(0.2), Inches(0.05), Inches(7.1))]:
            bar = sl.shapes.add_shape(1, *coords)
            bar.fill.solid(); bar.fill.fore_color.rgb = acc; bar.line.fill.background()
    except: pass

# Shablon uslubi → chizish funksiyasi xaritasi
STYLE_DRAW_MAP = {
    "minimalist_clean": draw_minimalist_clean,
    "minimalist_line": draw_minimalist_line,
    "minimalist_gray": draw_minimalist_clean,
    "minimalist_green": draw_minimalist_line,
    "minimalist_red": draw_minimalist_clean,
    "dark_elegant": draw_dark_elegant,
    "gradient_ocean": draw_gradient_ocean,
    "gradient_sunset": draw_gradient_ocean,
    "gradient_purple": draw_gradient_ocean,
    "gradient_forest": draw_gradient_ocean,
    "gradient_cosmic": draw_starnight,
    "gradient_peach": draw_gradient_ocean,
    "geometric_blue": draw_geometric_blue,
    "geometric_gold": draw_geometric_gold,
    "geometric_red": draw_geometric_blue,
    "geometric_mosaic": draw_geometric_mosaic,
    "geometric_bronze": draw_geometric_gold,
    "geometric_hex": draw_corporate_steel,
    "corporate_navy": draw_corporate_navy,
    "corporate_steel": draw_corporate_steel,
    "corporate_data": draw_minimalist_line,
    "corporate_target": draw_minimalist_clean,
    "corporate_classic": draw_corporate_navy,
    "corporate_academic": draw_minimalist_line,
    "creative_colorful": draw_geometric_blue,
    "creative_pink": draw_minimalist_clean,
    "creative_fire": draw_creative_fire,
    "creative_wave": draw_gradient_ocean,
    "creative_starnight": draw_starnight,
    "academic_classic": draw_academic,
}




# Bot uchun TEMPLATES alias
TEMPLATES = TEMPLATES_30


# ============================================================
# PPTX YARATISH
# ============================================================

def get_contrast_color(bg_rgb, text_rgb):
    """Fon va matn rangi bir xil bo'lsa, kontrastli rang qaytaradi"""
    def luminance(c):
        try:
            # RGBColor tuple kabi ishlaydi: c[0], c[1], c[2]
            return 0.299 * c[0] + 0.587 * c[1] + 0.114 * c[2]
        except Exception:
            return 128
    
    bg_lum = luminance(bg_rgb)
    txt_lum = luminance(text_rgb)
    
    diff = abs(bg_lum - txt_lum)
    if diff < 60:
        if bg_lum < 128:
            return RGBColor(255, 255, 255)
        else:
            return RGBColor(20, 20, 20)
    return text_rgb

def make_pptx(content, topic, tmpl_id, ud={}, user_imgs=None, img_pages=None):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    tmpl = TEMPLATES.get(str(tmpl_id), TEMPLATES["1"])
    bg1 = RGBColor(*tmpl["bg1"]); bg2 = RGBColor(*tmpl["bg2"])
    tc = RGBColor(*tmpl.get("title_c", tmpl.get("title", (255,255,255))))
    txc = RGBColor(*tmpl.get("text_c", tmpl.get("text", (220,220,220))))
    acc = RGBColor(*tmpl.get("accent", (255,200,0)))

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slaydlarni ajratish
    clean = content.replace("**","").replace("##","").replace("#","")
    slides_raw = []; cur_t = None; cur_b = []
    for line in clean.strip().split("\n"):
        line = line.strip()
        if not line: continue
        ul = line.upper()
        is_slide = (ul.startswith("SLAYD") or ul.startswith("SLIDE") or ul.startswith("СЛАЙД")) and ":" in line
        if is_slide:
            if cur_t is not None:
                slides_raw.append((cur_t, cur_b[:]))
            cur_t = line.split(":", 1)[1].strip()
            # Slayd raqamini sarlavhadan olib tashlash
            cur_t = re.sub(r"^\d+[.:]?\s*", "", cur_t).strip() or cur_t
            cur_b = []
        else:
            b = re.sub(r"^[-•►▸*\s]+", "", line)
            if b: cur_b.append(b)
    if cur_t is not None: slides_raw.append((cur_t, cur_b))
    # Agar SLAYD formati topilmasa, \n\n boyicha ajratish
    if not slides_raw or len(slides_raw) < 2:
        parts = [p.strip() for p in clean.split("\n\n") if p.strip()]
        slides_raw = []
        for p in parts:
            plines = [l.strip() for l in p.split("\n") if l.strip()]
            if plines:
                t_s = re.sub(r"^(SLAYD|SLIDE)\s*\d+[:\.]?\s*", "", plines[0], flags=re.IGNORECASE).strip()
                t_s = t_s or plines[0]
                b_s = plines[1:]
                slides_raw.append((t_s[:80], b_s))
    if not slides_raw:
        slides_raw = [(topic, [clean[:300]])]
    slides = slides_raw

    for sn, (title, bullets) in enumerate(slides):
        sl = prs.slides.add_slide(blank)

        # Gradient fon
        bg = sl.background; fill = bg.fill
        fill.gradient(); fill.gradient_angle = 2700000
        fill.gradient_stops[0].position = 0; fill.gradient_stops[0].color.rgb = bg1
        fill.gradient_stops[1].position = 1.0; fill.gradient_stops[1].color.rgb = bg2

        # Shablon uslubiga qarab dekorativ elementlar
        tmpl_style = TEMPLATES.get(str(tmpl_id), {}).get("style", "minimalist_clean")
        draw_func = STYLE_DRAW_MAP.get(tmpl_style, draw_minimalist_clean)
        try: draw_func(sl, TEMPLATES.get(str(tmpl_id), {}))
        except: pass

        # Yuqori chiziq
        try:
            bar = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.1))
            bar.fill.solid(); bar.fill.fore_color.rgb = acc; bar.line.fill.background()
        except: pass

        if sn == 0:
            # ── KUCHLI 1-SLAYD DIZAYNI ──
            # Chap panel
            try:
                lp = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(4.2), Inches(7.5))
                lp.fill.solid(); lp.fill.fore_color.rgb = acc
                lp.fill.fore_color.transparency = 0.2; lp.line.fill.background()
            except: pass
            # Diagonal accent
            try:
                da = sl.shapes.add_shape(1, Inches(4.0), Inches(0), Inches(0.18), Inches(7.5))
                da.fill.solid(); da.fill.fore_color.rgb = acc; da.line.fill.background()
            except: pass
            # O'ng yuqori kichik doira (dekorativ)
            try:
                dc = sl.shapes.add_shape(9, Inches(11.5), Inches(5.5), Inches(2.5), Inches(2.5))
                dc.fill.solid(); dc.fill.fore_color.rgb = acc
                dc.fill.fore_color.transparency = 0.85; dc.line.fill.background()
            except: pass
            # Katta sarlavha
            short_topic = topic if len(topic) <= 55 else topic[:52] + "..."
            tb = sl.shapes.add_textbox(Inches(4.5), Inches(1.0), Inches(8.5), Inches(3.0))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = short_topic
            p.font.size = Pt(36 if len(short_topic) < 35 else 28)
            p.font.bold = True; p.font.color.rgb = tc
            # Hook gap
            hook_map = {"uz": "Bilimni chuqurlashtiring, kelajakni quching",
                        "ru": "Углубите знания, постройте будущее",
                        "en": "Deepen knowledge, build the future"}
            hook = hook_map.get(ud.get("lang","uz"), hook_map["uz"])
            tb_h = sl.shapes.add_textbox(Inches(4.5), Inches(4.1), Inches(8.5), Inches(0.6))
            p_h = tb_h.text_frame.paragraphs[0]
            p_h.text = hook; p_h.font.size = Pt(15); p_h.font.italic = True
            p_h.font.color.rgb = txc
            # Separator
            try:
                sp = sl.shapes.add_shape(1, Inches(4.5), Inches(4.85), Inches(7.5), Inches(0.05))
                sp.fill.solid(); sp.fill.fore_color.rgb = acc; sp.line.fill.background()
            except: pass
            # Muallif ma'lumotlari
            info_lines = []
            if ud.get("full_name"): info_lines.append(f"Muallif: {ud['full_name']}")
            if ud.get("subject"): info_lines.append(f"Fan: {ud['subject']}")
            if ud.get("university"): info_lines.append(f"Universitet: {ud['university']}")
            if ud.get("teacher"): info_lines.append(f"O'qituvchi: {ud['teacher']}")
            info_lines.append(datetime.now().strftime("%Y-yil"))
            if info_lines:
                tb2 = sl.shapes.add_textbox(Inches(4.5), Inches(5.0), Inches(8.5), Inches(2.2))
                tf2 = tb2.text_frame; tf2.word_wrap = True; first2 = True
                for ln_txt in info_lines:
                    p2 = tf2.paragraphs[0] if first2 else tf2.add_paragraph(); first2 = False
                    p2.text = ln_txt; p2.font.size = Pt(14)
                    p2.font.color.rgb = txc; p2.space_before = Pt(3)
            # Chap panel matni
            try:
                tb_s = sl.shapes.add_textbox(Inches(0.2), Inches(2.5), Inches(3.7), Inches(2.5))
                tf_s = tb_s.text_frame; tf_s.word_wrap = True
                p_s = tf_s.paragraphs[0]
                p_s.text = (ud.get("subject") or topic)[:45]
                p_s.font.size = Pt(14); p_s.font.bold = True
                p_s.font.color.rgb = tc; p_s.alignment = PP_ALIGN.CENTER
            except: pass

        elif sn == 1 and any(w in title.upper() for w in ["REJA", "PLAN", "MUNDARIJA", "CONTENT"]):
            # 2-SLAYD: REJALAR
            tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(1.1))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = "R  E  J  A"
            p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = tc
            p.alignment = PP_ALIGN.CENTER
            try:
                ln2 = sl.shapes.add_shape(1, Inches(1.5), Inches(1.35), Inches(10.33), Inches(0.08))
                ln2.fill.solid(); ln2.fill.fore_color.rgb = acc; ln2.line.fill.background()
            except: pass
            if bullets:
                # Faqat bo'lim sarlavhalarini olish (uzun tavsiflarni emas)
                reja_items = []
                for b in bullets:
                    b = b.strip()
                    if not b: continue
                    # Faqat qisqa sarlavhalar (100 belgidan kam)
                    if len(b) < 120:
                        reja_items.append(b)
                    if len(reja_items) >= 8: break

                # Font o'lchami slayd sig'imiga qarab
                font_size = 22 if len(reja_items) <= 5 else (19 if len(reja_items) <= 7 else 16)
                spacing = 12 if len(reja_items) <= 5 else (8 if len(reja_items) <= 7 else 5)

                tb2 = sl.shapes.add_textbox(Inches(1.0), Inches(1.55), Inches(11.3), Inches(5.7))
                tf2 = tb2.text_frame; tf2.word_wrap = True; first = True
                num = 0
                for b in reja_items:
                    if not re.match(r'^[0-9]', b):
                        num += 1
                        display = f"{num}.   {b}"
                    else:
                        display = f"   {b}"
                    p2 = tf2.paragraphs[0] if first else tf2.add_paragraph(); first = False
                    p2.text = display
                    p2.font.size = Pt(font_size); p2.font.bold = False
                    p2.font.color.rgb = get_contrast_color(bg1, txc)
                    p2.space_before = Pt(spacing); p2.alignment = PP_ALIGN.LEFT

        else:
            # Oddiy slayd
            tb = sl.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.53), Inches(1.2))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = title[:80]
            p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = tc

            try:
                ln2 = sl.shapes.add_shape(1, Inches(0.4), Inches(1.52), Inches(12.53), Inches(0.07))
                ln2.fill.solid(); ln2.fill.fore_color.rgb = acc; ln2.line.fill.background()
            except: pass

            if bullets:
                has_img = (img_pages and any(img_pages.get(str(i)) == sn+1
                           for i in range(len(user_imgs or [])))) or \
                          (sn+1 in ud.get("ai_img_slides", []))
                txt_w = 7.8 if has_img else 12.5

                # Infografika va oddiy matn ajratish
                normal_b = []; infog = []
                for b in bullets:
                    if b.strip().upper().startswith("INFOGRAFIKA:"):
                        infog.append(b.strip())
                    elif b.strip():
                        normal_b.append(b.strip())

                tb2 = sl.shapes.add_textbox(Inches(0.4), Inches(1.75), Inches(txt_w), Inches(5.0))
                tf2 = tb2.text_frame; tf2.word_wrap = True; first = True
                # Barcha matnni birlashtirish - 250-350 so'z
                full_text = " ".join(normal_b)
                # Agar matn juda qisqa bo'lsa, barchasini ko'rsatish
                if len(normal_b) <= 3:
                    for b in normal_b:
                        p2 = tf2.paragraphs[0] if first else tf2.add_paragraph(); first = False
                        p2.text = b; p2.font.size = Pt(15); p2.font.color.rgb = txc
                        p2.space_before = Pt(6)
                else:
                    # Ko'p qatorli matn - paragraf sifatida ko'rsatish
                    for b in normal_b[:15]:
                        p2 = tf2.paragraphs[0] if first else tf2.add_paragraph(); first = False
                        # Uzun gaplar paragraf, qisqalar bullet
                        if len(b) > 80:
                            p2.text = b; p2.font.size = Pt(13)
                        else:
                            p2.text = f"▸  {b}"; p2.font.size = Pt(14)
                        p2.font.color.rgb = get_contrast_color(bg1, txc)
                        p2.space_before = Pt(3)

                # Diagramma qo'shish (faqat [DIAGRAMMA:] formati bo'lsa)
                full_slide_text = " ".join(bullets)
                # [DIAGRAMMA:...] matndan tozalash
                clean_bullets = []
                diag_data = parse_diagram_data(full_slide_text)
                for b in normal_b:
                    if not b.strip().upper().startswith("[DIAGRAMMA"):
                        clean_bullets.append(b)
                normal_b = clean_bullets
                
                with_diag = ud.get("with_diagram", False)
                if with_diag and diag_data and len(diag_data[0].get("data", [])) >= 2 and "[DIAGRAMMA" in " ".join(bullets).upper():
                    try:
                        add_diagram_to_slide(sl, topic, title, full_slide_text,
                                           diag_data, acc, get_contrast_color(bg1, txc), False)
                    except Exception as de:
                        logger.warning(f"Diagram error: {de}")

                # Infografika bar chart (eski)
                if infog:
                    try:
                        import re as _re2
                        for inf in infog[:1]:
                            parts = inf.split(":", 2)
                            inf_title = parts[1].strip() if len(parts) > 1 else ""
                            data_str = parts[2] if len(parts) > 2 else ""
                            entries = []
                            for item in data_str.split(","):
                                item = item.strip()
                                if ":" in item:
                                    k, v = item.split(":", 1)
                                    num = _re2.sub(r'[^\d.]', '', v.strip())
                                    try: entries.append((k.strip(), float(num), v.strip()))
                                    except: pass
                            if entries:
                                max_val = max(e[1] for e in entries) or 1
                                inf_y = 5.7
                                itb = sl.shapes.add_textbox(Inches(0.4), Inches(inf_y-0.35), Inches(12.0), Inches(0.3))
                                ip = itb.text_frame.paragraphs[0]
                                ip.text = f"📊 {inf_title}"; ip.font.size = Pt(12); ip.font.bold = True
                                ip.font.color.rgb = acc
                                for ei, (lbl, val, orig) in enumerate(entries[:5]):
                                    by = inf_y + ei * 0.35
                                    bw = min(7.5 * (val / max_val), 7.5)
                                    ltb = sl.shapes.add_textbox(Inches(0.4), Inches(by), Inches(2.4), Inches(0.28))
                                    ltb.text_frame.paragraphs[0].text = lbl[:18]
                                    ltb.text_frame.paragraphs[0].font.size = Pt(10)
                                    ltb.text_frame.paragraphs[0].font.color.rgb = txc
                                    bar = sl.shapes.add_shape(1, Inches(3.0), Inches(by+0.03), Inches(max(bw,0.2)), Inches(0.22))
                                    bar.fill.solid(); bar.fill.fore_color.rgb = acc; bar.line.fill.background()
                                    vtb = sl.shapes.add_textbox(Inches(3.0+bw+0.1), Inches(by), Inches(1.5), Inches(0.28))
                                    vtb.text_frame.paragraphs[0].text = orig
                                    vtb.text_frame.paragraphs[0].font.size = Pt(10)
                                    vtb.text_frame.paragraphs[0].font.bold = True
                                    vtb.text_frame.paragraphs[0].font.color.rgb = acc
                    except Exception as ie:
                        logger.error(f"Infografika: {ie}")

            # Sahifa raqami
            try:
                rq = sl.shapes.add_textbox(Inches(12.5), Inches(7.1), Inches(0.8), Inches(0.35))
                rq.text_frame.paragraphs[0].text = str(sn+1)
                rq.text_frame.paragraphs[0].font.size = Pt(11)
                rq.text_frame.paragraphs[0].font.color.rgb = acc
            except: pass

        # Foydalanuvchi rasmlari
        if user_imgs and img_pages:
            for ii, pn in img_pages.items():
                try:
                    ii_int = int(ii) if isinstance(ii, str) else ii
                    if ii_int < len(user_imgs) and pn == sn+1:
                        sl.shapes.add_picture(user_imgs[ii_int], Inches(8.8), Inches(1.6), Inches(4.2), Inches(5.4))
                except Exception as e: logger.error(f"User img: {e}")

    # AI rasmlar
    ai_slides = ud.get("ai_img_slides", [])
    if ai_slides:
        slide_list = list(prs.slides)
        for slide_num in ai_slides:
            try:
                if slide_num-1 < len(slide_list):
                    sl2 = slide_list[slide_num-1]
                    slide_title = topic
                    for sh in sl2.shapes:
                        if sh.has_text_frame and sh.text_frame.paragraphs:
                            txt = sh.text_frame.paragraphs[0].text.strip()
                            if txt and len(txt) > 3 and txt != topic:
                                slide_title = txt; break
                    img_buf = get_image(f"{slide_title} {topic}"[:60])
                    if img_buf:
                        if slide_num % 2 == 0:
                            sl2.shapes.add_picture(img_buf, Inches(0.4), Inches(1.7), Inches(12.5), Inches(3.2))
                        else:
                            sl2.shapes.add_picture(img_buf, Inches(8.6), Inches(1.7), Inches(4.5), Inches(3.5))
            except Exception as ie:
                logger.error(f"AI rasm {slide_num}: {ie}")

    td = tempfile.mkdtemp()
    out = os.path.join(td, "prezentatsiya.pptx")
    prs.save(out); return out, td

# ============================================================
# HTML PREZENTATSIYA
# ============================================================
def make_html(content, topic, tmpl_id, ud={}):
    tmpl = TEMPLATES.get(str(tmpl_id), TEMPLATES["1"])
    bg1 = "#{:02x}{:02x}{:02x}".format(*tmpl["bg1"])
    bg2 = "#{:02x}{:02x}{:02x}".format(*tmpl["bg2"])
    th = "#{:02x}{:02x}{:02x}".format(*tmpl.get("title_c", tmpl.get("title", (255,255,255))))
    tx = "#{:02x}{:02x}{:02x}".format(*tmpl.get("text_c", tmpl.get("text", (220,220,220))))
    ac = "#{:02x}{:02x}{:02x}".format(*tmpl.get("accent", (255,200,0)))

    clean = content.replace("**","").replace("##","").replace("#","")
    slides = []; cur_t = topic; cur_b = []
    for line in clean.strip().split("\n"):
        line = line.strip()
        if not line: continue
        ul = line.upper()
        is_slide = (ul.startswith("SLAYD") or ul.startswith("SLIDE")) and ":" in line
        if is_slide:
            if cur_t is not None: slides.append((cur_t, cur_b[:]))
            cur_t = line.split(":", 1)[1].strip(); cur_b = []
        else:
            b = re.sub(r'^[-•►▸*\s]+', '', line)
            if b: cur_b.append(b)
    if cur_t: slides.append((cur_t, cur_b))
    total = len(slides)
    info = build_info(ud)

    slides_html = ""
    for i, (ti, b) in enumerate(slides):
        bl = "".join(f"<li>{x}</li>" for x in b[:10] if x.strip())
        disp = "block" if i == 0 else "none"
        if i == 0:
            slides_html += f"""<div class="slide" id="s{i}" style="display:{disp}">
<div class="snum">{i+1}/{total}</div>
<div class="title-slide"><h1>{ti}</h1>
<div class="info">{info.replace(chr(10),"<br>")}</div>
<p class="yr">{datetime.now().strftime("%Y-yil")}</p></div></div>"""
        else:
            slides_html += f"""<div class="slide" id="s{i}" style="display:{disp}">
<div class="snum">{i+1}/{total}</div>
<h2>{ti}</h2><div class="aline"></div>
<ul>{bl}</ul></div>"""

    html = f"""<!DOCTYPE html><html lang="uz"><head><meta charset="UTF-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>{topic}</title><style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{font-family:'Segoe UI',Arial,sans-serif;background:#0a0a0a;display:flex;justify-content:center;align-items:center;min-height:100vh;flex-direction:column;gap:20px;padding:20px}}
.slide{{background:linear-gradient(135deg,{bg1},{bg2});width:900px;max-width:98vw;aspect-ratio:16/9;border-radius:14px;padding:45px 55px;position:relative;overflow:hidden;box-shadow:0 25px 70px rgba(0,0,0,.6)}}
.slide::before{{content:'';position:absolute;right:-70px;top:-70px;width:280px;height:280px;border-radius:50%;background:rgba(255,255,255,.06)}}
.slide::after{{content:'';position:absolute;left:-50px;bottom:-50px;width:200px;height:200px;border-radius:50%;background:rgba(255,255,255,.04)}}
.snum{{position:absolute;bottom:18px;right:25px;color:{ac};font-size:14px;font-weight:600}}
.title-slide{{display:flex;flex-direction:column;justify-content:center;align-items:center;height:100%;text-align:center}}
.title-slide h1{{color:{th};font-size:40px;font-weight:800;line-height:1.25;margin-bottom:20px;text-shadow:0 2px 10px rgba(0,0,0,.3)}}
.info{{color:{tx};font-size:17px;line-height:2;opacity:.92;margin-bottom:12px}}
.yr{{color:{ac};font-size:16px;font-weight:700}}
h2{{color:{th};font-size:26px;font-weight:700;margin-bottom:6px}}
.aline{{width:60px;height:4px;background:{ac};border-radius:2px;margin:10px 0 20px}}
ul{{list-style:none;color:{tx}}}
li{{font-size:18px;line-height:1.7;margin-bottom:10px;padding-left:24px;position:relative}}
li::before{{content:'▸';position:absolute;left:0;color:{ac};font-weight:bold}}
.ctrl{{display:flex;gap:12px;align-items:center;margin-top:5px}}
.btn{{background:rgba(255,255,255,.12);color:white;border:2px solid {ac};padding:10px 30px;border-radius:30px;cursor:pointer;font-size:15px;font-weight:600;transition:all .3s;backdrop-filter:blur(10px)}}
.btn:hover{{background:{ac};color:#000;transform:scale(1.05)}}
.prog{{color:rgba(255,255,255,.6);font-size:14px;min-width:60px;text-align:center}}
</style></head><body>
{slides_html}
<div class="ctrl">
<button class="btn" onclick="nav(-1)">◀ Oldingi</button>
<span class="prog" id="pg">1/{total}</span>
<button class="btn" onclick="nav(1)">Keyingi ▶</button>
</div>
<script>
var cur=0,tot={total};
function nav(d){{
var ns=cur+d;
if(ns<0||ns>=tot)return;
document.getElementById('s'+cur).style.display='none';
cur=ns;
document.getElementById('s'+cur).style.display='block';
document.getElementById('pg').textContent=(cur+1)+'/'+tot;
}}
document.addEventListener('keydown',function(e){{
if(e.key==='ArrowRight'||e.key===' ')nav(1);
if(e.key==='ArrowLeft')nav(-1);
}});
</script></body></html>"""

    td = tempfile.mkdtemp()
    out = os.path.join(td, "prezentatsiya.html")
    with open(out, "w", encoding="utf-8") as f: f.write(html)
    return out, td

# ============================================================
# DOCX YARATISH
# ============================================================
def make_docx(content, topic, ud={}):
    try:
        from docx import Document
        from docx.shared import Pt as DPt, Inches as DInches, RGBColor as DRGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()
        style = doc.styles['Normal']
        style.font.name = 'Times New Roman'
        style.font.size = DPt(14)

        # Sarlavha sahifasi
        if ud.get("university"):
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run(ud["university"].upper())
            run.bold = True; run.font.size = DPt(14)

        doc.add_paragraph()
        title_p = doc.add_heading(topic, 0)
        title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER

        doc.add_paragraph()
        if ud.get("full_name"):
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            p.add_run(f"Tayyorladi: {ud['full_name']}")
        if ud.get("teacher"):
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            p.add_run(f"O'qituvchi: {ud['teacher']}")
        if ud.get("city"):
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            p.add_run(f"{ud.get('city','')} — {datetime.now().year}")

        doc.add_page_break()

        # Kontent
        for line in content.split("\n"):
            line = line.strip()
            if not line: continue
            if any(line.startswith(h) for h in ["KIRISH","XULOSA","I BOB","II BOB","III BOB",
                   "MUNDARIJA","ANNOTATSIYA","ABSTRACT","METODOLOGIYA","NATIJALAR",
                   "ADABIYOTLAR","FOYDALANILGAN"]):
                h = doc.add_heading(line, 1)
                h.alignment = WD_ALIGN_PARAGRAPH.CENTER
            else:
                p = doc.add_paragraph(line)
                p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                for run in p.runs:
                    run.font.size = DPt(14)

        td = tempfile.mkdtemp()
        out = os.path.join(td, "dokument.docx")
        doc.save(out); return out, td
    except Exception as e:
        logger.error(f"DOCX xato: {e}"); return None, None

# ============================================================
# PDF YARATISH
# ============================================================
def make_pdf(content, topic, ud={}):
    try:
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import cm
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        import os as _os

        font_paths = [
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
        ]
        font_name = "Helvetica"
        for fp in font_paths:
            if _os.path.exists(fp):
                try:
                    pdfmetrics.registerFont(TTFont("CustomFont", fp))
                    font_name = "CustomFont"; break
                except: pass

        td = tempfile.mkdtemp()
        out = os.path.join(td, "dokument.pdf")
        doc = SimpleDocTemplate(out, pagesize=A4,
                                rightMargin=2*cm, leftMargin=2*cm,
                                topMargin=2*cm, bottomMargin=2*cm)
        styles = getSampleStyleSheet()
        story = []

        # Sarlavha
        title_style = ParagraphStyle("title", fontName=font_name, fontSize=16,
                                     spaceAfter=20, alignment=1, leading=20)
        story.append(Paragraph(topic, title_style))
        story.append(Spacer(1, 20))

        # Muallif
        info_style = ParagraphStyle("info", fontName=font_name, fontSize=12,
                                    spaceAfter=6, alignment=1)
        info = build_info(ud)
        for line in info.split("\n"):
            if line.strip():
                story.append(Paragraph(line, info_style))
        story.append(Spacer(1, 30))

        # Kontent
        head_style = ParagraphStyle("head", fontName=font_name, fontSize=14,
                                    spaceAfter=10, spaceBefore=15, fontWeight=700, leading=18)
        body_style = ParagraphStyle("body", fontName=font_name, fontSize=12,
                                    spaceAfter=8, leading=18, alignment=4)

        for line in content.split("\n"):
            line = line.strip()
            if not line: story.append(Spacer(1, 6)); continue
            if any(line.startswith(h) for h in ["KIRISH","XULOSA","I BOB","II BOB","III BOB",
                   "MUNDARIJA","ANNOTATSIYA","METODOLOGIYA","NATIJALAR","ADABIYOTLAR","FOYDALANILGAN"]):
                story.append(Paragraph(line, head_style))
            else:
                story.append(Paragraph(line.replace("<","&lt;").replace(">","&gt;"), body_style))

        doc.build(story)
        return out, td
    except Exception as e:
        logger.error(f"PDF xato: {e}"); return None, None

# ============================================================
# KONVERTATSIYA
# ============================================================
def pdf_to_text(pdf_path):
    try:
        import fitz
        d = fitz.open(pdf_path)
        return " ".join(pg.get_text() for pg in d)
    except: return ""

def imgs_to_pdf(imgs, out_path):
    try:
        from PIL import Image
        images = []
        for p in imgs:
            img = Image.open(p).convert("RGB")
            images.append(img)
        if images:
            images[0].save(out_path, save_all=True, append_images=images[1:])
            return True
    except Exception as e:
        logger.error(f"imgs_to_pdf: {e}")
    return False

# ============================================================
# TUGMALAR
# ============================================================
def main_kb(uid):
    kb = types.ReplyKeyboardMarkup(resize_keyboard=True)
    kb.row(t(uid,"btn_referat"), t(uid,"btn_kurs"))
    kb.row(t(uid,"btn_mustaqil"), t(uid,"btn_maqola"))
    kb.row(t(uid,"btn_prez"), t(uid,"btn_test"))
    kb.row(t(uid,"btn_imlo"), t(uid,"btn_konv"))
    kb.row(t(uid,"btn_balans"), t(uid,"btn_orders"))
    kb.row(t(uid,"btn_referral"), t(uid,"btn_donat"))
    kb.row(t(uid,"btn_help"), t(uid,"btn_admin"))
    return kb

def lang_kb():
    kb = types.InlineKeyboardMarkup(row_width=3)
    kb.add(
        types.InlineKeyboardButton("🇺🇿 O'zbek", callback_data="lang:uz"),
        types.InlineKeyboardButton("🇷🇺 Русский", callback_data="lang:ru"),
        types.InlineKeyboardButton("🇬🇧 English", callback_data="lang:en")
    )
    return kb

def bk_kb():
    kb = types.InlineKeyboardMarkup(row_width=2)
    kb.add(
        types.InlineKeyboardButton("◀️ Orqaga", callback_data="back_step"),
        types.InlineKeyboardButton("🏠 Menyu", callback_data="bk")
    )
    return kb

def skip_kb(ns):
    kb = types.InlineKeyboardMarkup(row_width=2)
    kb.add(
        types.InlineKeyboardButton("⏭ O'tkazib yuborish", callback_data=f"skip:{ns}"),
        types.InlineKeyboardButton("🏠 Menyu", callback_data="bk")
    )
    return kb

def fmt_kb(prefix):
    kb = types.InlineKeyboardMarkup(row_width=3)
    kb.add(
        types.InlineKeyboardButton("📝 DOCX", callback_data=f"{prefix}:docx"),
        types.InlineKeyboardButton("📄 PDF", callback_data=f"{prefix}:pdf"),
        types.InlineKeyboardButton("📱 TXT", callback_data=f"{prefix}:txt")
    )
    return kb

def prez_fmt_kb():
    kb = types.InlineKeyboardMarkup(row_width=1)
    kb.add(
        types.InlineKeyboardButton("📊 PPTX (PowerPoint)", callback_data="pfmt:pptx"),
        types.InlineKeyboardButton("🌐 HTML (Interaktiv)", callback_data="pfmt:html"),
        types.InlineKeyboardButton("📦 Ikkalasi ham", callback_data="pfmt:both")
    )
    return kb

def slides_kb():
    kb = types.InlineKeyboardMarkup(row_width=4)
    for n in [10,15,20,25,30,35,40,50]:
        kb.add(types.InlineKeyboardButton(f"{n} slayd — {n*PRICE_SLIDE:,} so'm", callback_data=f"slides:{n}"))
    kb.add(types.InlineKeyboardButton("✏️ O'zim yozaman", callback_data="slides:custom"))
    kb.add(types.InlineKeyboardButton("🏠 Menyu", callback_data="bk"))
    return kb

def test_kb():
    kb = types.InlineKeyboardMarkup(row_width=4)
    for n in [10,20,30,50,100,200,500,1000]:
        kb.add(types.InlineKeyboardButton(f"{n} ta — {n*PRICE_TEST:,} so'm", callback_data=f"tcount:{n}"))
    kb.add(types.InlineKeyboardButton("✏️ O'zim yozaman", callback_data="tcount:custom"))
    kb.add(types.InlineKeyboardButton("🏠 Menyu", callback_data="bk"))
    return kb

def plans_kb():
    kb = types.InlineKeyboardMarkup(row_width=4)
    for n in [3,4,5,6,7,8,10,12]:
        kb.add(types.InlineKeyboardButton(f"{n} ta bo'lim", callback_data=f"plans:{n}"))
    kb.add(types.InlineKeyboardButton("✏️ O'zim yozaman", callback_data="plans:custom"))
    kb.add(types.InlineKeyboardButton("🏠 Menyu", callback_data="bk"))
    return kb

def lc_kb(prefix):
    kb = types.InlineKeyboardMarkup(row_width=3)
    kb.add(
        types.InlineKeyboardButton("🇺🇿 O'zbek", callback_data=f"{prefix}:uz"),
        types.InlineKeyboardButton("🇷🇺 Rus", callback_data=f"{prefix}:ru"),
        types.InlineKeyboardButton("🇬🇧 Ingliz", callback_data=f"{prefix}:en")
    )
    return kb

def tmpl_kb(page=0):
    kb = types.InlineKeyboardMarkup(row_width=2)
    keys = list(TEMPLATES.keys())
    per = 6; start = page * per; end = min(start + per, len(keys))
    for k in keys[start:end]:
        tmpl = TEMPLATES[k]
        preview = tmpl.get("preview", "🎨")
        name = tmpl.get("name", k)
        desc = tmpl.get("desc", "")
        btn_text = f"{name}"
        kb.add(types.InlineKeyboardButton(btn_text, callback_data=f"tmpl:{k}"))
    nav_btns = []
    if page > 0: nav_btns.append(types.InlineKeyboardButton("◀️ Oldingi", callback_data=f"tmpl_p:{page-1}"))
    if end < len(keys): nav_btns.append(types.InlineKeyboardButton("Keyingi ▶️", callback_data=f"tmpl_p:{page+1}"))
    if nav_btns: kb.row(*nav_btns)
    # Sahifa ko'rsatkichi
    total_pages = (len(keys) + per - 1) // per
    kb.add(types.InlineKeyboardButton(f"📄 {page+1}/{total_pages} sahifa", callback_data="noop"))
    return kb


def source_kb():
    """Manba tanlash — kitob yoki mavzu"""
    kb = types.InlineKeyboardMarkup(row_width=1)
    kb.add(
        types.InlineKeyboardButton("📚 Kitob/Fan nomi yozaman", callback_data="src:text"),
        types.InlineKeyboardButton("📄 PDF kitob yuklayaman", callback_data="src:pdf"),
        types.InlineKeyboardButton("🌐 Umumiy mavzu (manbasisiz)", callback_data="src:none"),
        types.InlineKeyboardButton("🏠 Menyu", callback_data="bk")
    )
    return kb

def plans_kb():
    """Nechta bo'lim bo'lsin tugmalari"""
    kb = types.InlineKeyboardMarkup(row_width=3)
    kb.add(
        types.InlineKeyboardButton("3 ta", callback_data="plans:3"),
        types.InlineKeyboardButton("4 ta", callback_data="plans:4"),
        types.InlineKeyboardButton("5 ta", callback_data="plans:5"),
        types.InlineKeyboardButton("6 ta", callback_data="plans:6"),
        types.InlineKeyboardButton("7 ta", callback_data="plans:7"),
        types.InlineKeyboardButton("8 ta", callback_data="plans:8"),
    )
    kb.add(types.InlineKeyboardButton("🏠 Menyu", callback_data="bk"))
    return kb

def img_choice_kb():
    kb = types.InlineKeyboardMarkup(row_width=1)
    kb.add(
        types.InlineKeyboardButton("🤖 AI avtomatik rasm qo'ysin", callback_data="img:ai"),
        types.InlineKeyboardButton("🖼 O'zim rasm yuklayman", callback_data="img:user"),
        types.InlineKeyboardButton("❌ Rasmsiz davom etish", callback_data="img:none")
    )
    return kb

def img_slide_select_kb(total_slides, page=0, mode="ai"):
    """Qaysi slaydlarga rasm yuklashni tanlash — sahifalab"""
    per_page = 10
    start = page * per_page + 1
    end = min(start + per_page - 1, total_slides)
    kb = types.InlineKeyboardMarkup(row_width=5)
    btns = []
    for n in range(start, end + 1):
        btns.append(types.InlineKeyboardButton(str(n), callback_data=f"img_slide:{n}:{mode}"))
    kb.add(*btns)
    nav = []
    if page > 0:
        nav.append(types.InlineKeyboardButton("◀️", callback_data=f"img_slide_page:{page-1}:{mode}"))
    if end < total_slides:
        nav.append(types.InlineKeyboardButton("▶️", callback_data=f"img_slide_page:{page+1}:{mode}"))
    if nav: kb.row(*nav)
    total_pages = (total_slides + per_page - 1) // per_page
    kb.add(types.InlineKeyboardButton(f"📄 {page+1}/{total_pages} ({start}-{end})", callback_data="noop"))
    kb.add(types.InlineKeyboardButton("✅ Tayyor (tanlangan slaydlar bilan davom et)", callback_data=f"img_slide_done:{mode}"))
    kb.add(types.InlineKeyboardButton("🏠 Menyu", callback_data="bk"))
    return kb

def conv_kb():
    kb = types.InlineKeyboardMarkup()
    kb.add(types.InlineKeyboardButton("📄➡️📊 PDF → PPTX", callback_data="cv:pdf"))
    kb.add(types.InlineKeyboardButton("📊➡️📄 PPTX → PDF", callback_data="cv:pptx"))
    kb.add(types.InlineKeyboardButton("🖼➡️📄 Rasmlar → PDF", callback_data="cv:img"))
    kb.add(types.InlineKeyboardButton("🔙 Orqaga", callback_data="bk"))
    return kb

# ============================================================
# INFO STEPS
# ============================================================
INFO_STEPS = [
    ("ask_name",    "full_name",   True),
    ("ask_univ",    "university",  True),
    ("ask_faculty", "faculty",     False),
    ("ask_year",    "year",        False),
    ("ask_teacher", "teacher",     False),
    ("ask_subject", "subject",     False),
    ("ask_city",    "city",        False),
]
INFO_STATES = [s[0] for s in INFO_STEPS]

def finish_info(uid, ud):
    svc = ud.get("svc", "referat")
    if svc == "prez":
        sst(uid, "prez_tmpl")
        bot.send_message(uid, t(uid, "ask_template"), reply_markup=tmpl_kb())
    else:
        sst(uid, f"{svc}_lang")
        bot.send_message(uid, t(uid, "ask_lang"), reply_markup=lc_kb(f"{svc}_lang"))

# ============================================================
# BUYRUQLAR
# ============================================================

# ============================================================
# REFERAL TIZIMI
# ============================================================
REFERAL_BONUS = 500  # Har bir do'st uchun bonus

def save_referral(referrer_id, referred_id):
    """Referal bog'liqligini saqlash"""
    try:
        c = sqlite3.connect("edubot.db"); cur = c.cursor()
        cur.execute("SELECT id FROM referrals WHERE referred_id=?", (referred_id,))
        if cur.fetchone():
            c.close(); return False
        cur.execute(
            "INSERT INTO referrals(referrer_id, referred_id, bonus_paid, created_at) VALUES(?,?,0,?)",
            (referrer_id, referred_id, datetime.now().strftime("%d.%m.%Y %H:%M")))
        cur.execute("UPDATE users SET referred_by=? WHERE telegram_id=?",
            (referrer_id, referred_id))
        c.commit(); c.close(); return True
    except Exception as e:
        logger.error(f"save_referral: {e}"); return False

def pay_referral_bonus(referred_id):
    """Do'sti birinchi buyurtma qilganda bonus to'lash"""
    try:
        c = sqlite3.connect("edubot.db"); cur = c.cursor()
        cur.execute(
            "SELECT id, referrer_id FROM referrals WHERE referred_id=? AND bonus_paid=0",
            (referred_id,))
        row = cur.fetchone()
        if not row:
            c.close(); return
        ref_id, referrer_id = row
        # Bonus berish
        cur.execute("UPDATE users SET balance=balance+? WHERE telegram_id=?",
            (REFERAL_BONUS, referrer_id))
        cur.execute("UPDATE referrals SET bonus_paid=1 WHERE id=?", (ref_id,))
        c.commit(); c.close()
        # Taklif qilganga xabar
        try:
            ref_user = get_user(referred_id)
            fname = ref_user["first_name"] if ref_user else "Do'stingiz"
            bot.send_message(referrer_id,
                f"🎉 *{fname}* sizning taklifingiz orqali birinchi buyurtma berdi!\n"
                f"💰 *{REFERAL_BONUS:,} so'm* bonus hisobingizga qo'shildi!\n"
                f"💳 Joriy balans: *{get_balance(referrer_id):,} so'm*",
                parse_mode="Markdown")
        except: pass
        logger.info(f"Referral bonus paid: {referrer_id} <- {referred_id}")
    except Exception as e:
        logger.error(f"pay_referral_bonus: {e}")

def get_referral_stats(uid):
    """Foydalanuvchining referal statistikasi"""
    try:
        c = sqlite3.connect("edubot.db"); cur = c.cursor()
        cur.execute("SELECT COUNT(*) FROM referrals WHERE referrer_id=?", (uid,))
        total = cur.fetchone()[0]
        cur.execute("SELECT COUNT(*) FROM referrals WHERE referrer_id=? AND bonus_paid=1", (uid,))
        paid = cur.fetchone()[0]
        cur.execute("SELECT COUNT(*) FROM referrals WHERE referrer_id=? AND bonus_paid=0", (uid,))
        pending = cur.fetchone()[0]
        c.close()
        return total, paid, pending
    except: return 0, 0, 0

def get_referral_link(uid):
    """Foydalanuvchi uchun referal havola"""
    bot_info = bot.get_me()
    return f"https://t.me/{bot_info.username}?start=ref_{uid}"

# ============================================================
# BALANS TO'LDIRISH TIZIMI
# ============================================================
TOPUP_AMOUNTS = [5000, 10000, 15000, 20000, 25000, 30000, 35000, 40000, 45000, 50000]

def save_topup_request(uid, amount):
    """To'ldirish so'rovini saqlash"""
    try:
        c = sqlite3.connect("edubot.db"); cur = c.cursor()
        cur.execute(
            "INSERT INTO topup_requests(telegram_id, amount, status, created_at) VALUES(?,?,'pending',?)",
            (uid, amount, datetime.now().strftime("%d.%m.%Y %H:%M")))
        c.commit()
        req_id = cur.lastrowid
        c.close(); return req_id
    except Exception as e:
        logger.error(f"save_topup: {e}"); return None

def approve_topup(req_id):
    """To'ldirish so'rovini tasdiqlash"""
    try:
        c = sqlite3.connect("edubot.db"); cur = c.cursor()
        cur.execute("SELECT telegram_id, amount FROM topup_requests WHERE id=? AND status='pending'",
            (req_id,))
        row = cur.fetchone()
        if not row:
            c.close(); return None, None
        uid, amount = row
        cur.execute("UPDATE topup_requests SET status='approved' WHERE id=?", (req_id,))
        cur.execute("UPDATE users SET balance=balance+? WHERE telegram_id=?", (amount, uid))
        c.commit(); c.close()
        return uid, amount
    except Exception as e:
        logger.error(f"approve_topup: {e}"); return None, None

def reject_topup(req_id):
    """To'ldirish so'rovini rad etish"""
    try:
        c = sqlite3.connect("edubot.db"); cur = c.cursor()
        cur.execute("SELECT telegram_id FROM topup_requests WHERE id=?", (req_id,))
        row = cur.fetchone()
        cur.execute("UPDATE topup_requests SET status='rejected' WHERE id=?", (req_id,))
        c.commit(); c.close()
        return row[0] if row else None
    except: return None

def topup_kb():
    """Balans to'ldirish miqdor tugmalari"""
    kb = types.InlineKeyboardMarkup(row_width=2)
    btns = []
    for amt in TOPUP_AMOUNTS:
        btns.append(types.InlineKeyboardButton(
            f"{amt:,} so'm", callback_data=f"topup_amt:{amt}"))
    kb.add(*btns)
    kb.add(types.InlineKeyboardButton("🏠 Menyu", callback_data="bk"))
    return kb

def referral_kb(uid):
    """Referal menyu tugmalari"""
    kb = types.InlineKeyboardMarkup(row_width=1)
    kb.add(
        types.InlineKeyboardButton("🔗 Referal havolam", callback_data="ref:link"),
        types.InlineKeyboardButton("📊 Statistika", callback_data="ref:stats"),
        types.InlineKeyboardButton("🏠 Menyu", callback_data="bk")
    )
    return kb

@bot.message_handler(commands=["start"])
def cmd_start(msg):
    uid = msg.from_user.id
    uname = msg.from_user.username or ""
    fname = msg.from_user.first_name or ""
    is_new = reg_user(uid, uname, fname)

    # Referal parametrini tekshirish
    args = msg.text.split() if msg.text else []
    if len(args) > 1 and args[1].startswith("ref_") and is_new:
        try:
            referrer_id = int(args[1][4:])
            if referrer_id != uid:
                save_referral(referrer_id, uid)
        except: pass

    # Obuna tekshirish
    if not check_subscription(uid):
        channels = get_sub_channels()
        ch_list = "\n".join([f"• {name}" for _, name in channels])
        bot.send_message(uid,
            f"⚠️ *Botdan foydalanish uchun quyidagi kanallarga obuna bo\'ling:*\n\n{ch_list}\n\n"
            f"Obuna bo\'lgach ✅ tugmasini bosing.",
            parse_mode="Markdown", reply_markup=sub_check_kb())
        return

    txt = t(uid, "welcome", name=fname)
    if is_new:
        txt += t(uid, "bonus", amount=BONUS_FIRST)
        if len(args) > 1 and args[1].startswith("ref_"):
            txt += "\n\n🎁 Siz taklif orqali keldingiz! Birinchi buyurtmadan so'ng taklif qilgan do'stingiz bonus oladi."
    txt += t(uid, "choose_lang")
    bot.send_message(uid, txt, parse_mode="Markdown", reply_markup=lang_kb())
    if is_new:
        try: bot.send_message(ADMIN_ID, f"🆕 Yangi: {fname} (@{uname}) | ID: {uid}")
        except: pass

@bot.message_handler(commands=["referat"])
def cmd_referat(msg):
    uid = msg.from_user.id
    reg_user(uid, msg.from_user.username or "", msg.from_user.first_name or "")
    if not check_subscription(uid):
        bot.send_message(uid, "⚠️ Botdan foydalanish uchun kanallarga obuna bo\'ling!", reply_markup=sub_check_kb())
        return
    UD.setdefault(uid, {})["source_type"] = "none"
    sst(uid, "referat_t", svc="referat")
    bot.send_message(uid, t(uid, "enter_topic"), reply_markup=bk_kb())

@bot.message_handler(commands=["kursishi"])
def cmd_kurs(msg):
    uid = msg.from_user.id
    reg_user(uid, msg.from_user.username or "", msg.from_user.first_name or "")
    if not check_subscription(uid):
        bot.send_message(uid, "⚠️ Botdan foydalanish uchun kanallarga obuna bo\'ling!", reply_markup=sub_check_kb())
        return
    sst(uid, "kurs_t", svc="kurs")
    bot.send_message(uid, "📚 Qaysi manbadan foydalanaylik?", reply_markup=source_kb())

@bot.message_handler(commands=["mustaqilish"])
def cmd_mustaqil(msg):
    uid = msg.from_user.id
    reg_user(uid, msg.from_user.username or "", msg.from_user.first_name or "")
    if not check_subscription(uid):
        bot.send_message(uid, "⚠️ Botdan foydalanish uchun kanallarga obuna bo\'ling!", reply_markup=sub_check_kb())
        return
    sst(uid, "mustaqil_t", svc="mustaqil")
    bot.send_message(uid, "📚 Qaysi manbadan foydalanaylik?", reply_markup=source_kb())

@bot.message_handler(commands=["maqola"])
def cmd_maqola(msg):
    uid = msg.from_user.id
    reg_user(uid, msg.from_user.username or "", msg.from_user.first_name or "")
    if not check_subscription(uid):
        bot.send_message(uid, "⚠️ Botdan foydalanish uchun kanallarga obuna bo\'ling!", reply_markup=sub_check_kb())
        return
    sst(uid, "maqola_t", svc="maqola")
    bot.send_message(uid, "📚 Qaysi manbadan foydalanaylik?", reply_markup=source_kb())

@bot.message_handler(commands=["prezentatsiya"])
def cmd_prez(msg):
    uid = msg.from_user.id
    reg_user(uid, msg.from_user.username or "", msg.from_user.first_name or "")
    sst(uid, "prez_t", svc="prez")
    bot.send_message(uid, t(uid, "enter_topic"), reply_markup=bk_kb())

@bot.message_handler(commands=["test"])
def cmd_test(msg):
    uid = msg.from_user.id
    reg_user(uid, msg.from_user.username or "", msg.from_user.first_name or "")
    sst(uid, "test_t", svc="test")
    bot.send_message(uid, t(uid, "enter_topic"), reply_markup=bk_kb())

@bot.message_handler(commands=["imlo"])
def cmd_imlo(msg):
    uid = msg.from_user.id
    reg_user(uid, msg.from_user.username or "", msg.from_user.first_name or "")
    sst(uid, "imlo_t")
    kb2 = types.InlineKeyboardMarkup()
    kb2.add(types.InlineKeyboardButton("📁 Fayl yuborish (PDF/TXT)", callback_data="imlo_file"))
    bot.send_message(uid, t(uid, "imlo_prompt"), reply_markup=kb2)

@bot.message_handler(commands=["konvertatsiya"])
def cmd_konv(msg):
    uid = msg.from_user.id
    reg_user(uid, msg.from_user.username or "", msg.from_user.first_name or "")
    bot.send_message(uid, "🔄 Format tanlang:", reply_markup=conv_kb())

@bot.message_handler(commands=["balans"])
def cmd_balans(msg):
    uid = msg.from_user.id
    reg_user(uid, msg.from_user.username or "", msg.from_user.first_name or "")
    bal = get_balance(uid)
    kb2 = types.InlineKeyboardMarkup()
    kb2.add(types.InlineKeyboardButton(t(uid,"topup"), callback_data="topup"))
    bot.send_message(uid, t(uid, "balance_info", bal=bal), parse_mode="Markdown", reply_markup=kb2)

@bot.message_handler(commands=["menu"])
def cmd_menu(msg):
    uid = msg.from_user.id
    reg_user(uid, msg.from_user.username or "", msg.from_user.first_name or "")
    bot.send_message(uid, t(uid, "menu"), reply_markup=main_kb(uid))

@bot.message_handler(commands=["yordam", "help"])
def cmd_help(msg):
    uid = msg.from_user.id
    help_txt = t(uid, "help_text") + "\n\n" + t(uid, "prices",
        p1=PRICE_PAGE, p2=PRICE_KURS, p3=PRICE_MUSTAQIL,
        p4=PRICE_MAQOLA, p5=PRICE_SLIDE, p6=PRICE_TEST)
    bot.send_message(uid, help_txt, parse_mode="Markdown", reply_markup=main_kb(uid))


@bot.message_handler(commands=["subon"])
def cmd_subon(msg):
    if msg.from_user.id != ADMIN_ID: return
    global SUB_ENABLED
    SUB_ENABLED = True
    bot.send_message(msg.chat.id, "✅ Majburiy obuna *yoqildi!*", parse_mode="Markdown")

@bot.message_handler(commands=["suboff"])
def cmd_suboff(msg):
    if msg.from_user.id != ADMIN_ID: return
    global SUB_ENABLED
    SUB_ENABLED = False
    bot.send_message(msg.chat.id, "❌ Majburiy obuna *o'chirildi!*", parse_mode="Markdown")

@bot.message_handler(commands=["addchannel"])
def cmd_addchannel(msg):
    if msg.from_user.id != ADMIN_ID: return
    try:
        parts = msg.text.split(None, 2)
        ch_id = parts[1]
        ch_name = parts[2] if len(parts) > 2 else ch_id
        add_sub_channel(ch_id, ch_name)
        bot.send_message(msg.chat.id, f"✅ Kanal qo'shildi: *{ch_name}*", parse_mode="Markdown")
    except:
        bot.send_message(msg.chat.id, "❌ /addchannel [kanal_id] [kanal_nomi]\nMasalan: /addchannel -1001234567890 Mening Kanalim")

@bot.message_handler(commands=["removechannel"])
def cmd_removechannel(msg):
    if msg.from_user.id != ADMIN_ID: return
    try:
        ch_id = msg.text.split()[1]
        remove_sub_channel(ch_id)
        bot.send_message(msg.chat.id, f"✅ Kanal o'chirildi: {ch_id}")
    except:
        bot.send_message(msg.chat.id, "❌ /removechannel [kanal_id]")

@bot.message_handler(commands=["channels"])
def cmd_channels(msg):
    if msg.from_user.id != ADMIN_ID: return
    channels = get_sub_channels()
    if not channels:
        bot.send_message(msg.chat.id, "📋 Hech qanday kanal qo'shilmagan.")
        return
    txt = f"📋 *Kanallar ({len(channels)} ta):*\n\n"
    for ch_id, ch_name in channels:
        txt += f"• {ch_name} | ID: `{ch_id}`\n"
    txt += f"\n{'✅ Obuna yoqilgan' if SUB_ENABLED else '❌ Obuna o\'chirilgan'}"
    bot.send_message(msg.chat.id, txt, parse_mode="Markdown")

@bot.message_handler(commands=["stats"])
def cmd_stats(msg):
    if msg.from_user.id != ADMIN_ID: return
    u, w, c, i = get_stats()
    bot.send_message(msg.chat.id,
        f"📊 *Statistika*\n\n👥 Foydalanuvchilar: {u}\n"
        f"📝 Ishlar: {w}\n🔄 Konvertatsiyalar: {c}\n💰 Daromad: {i:,} so'm",
        parse_mode="Markdown")

@bot.message_handler(commands=["broadcast"])
def cmd_bc(msg):
    if msg.from_user.id != ADMIN_ID: return
    sst(msg.from_user.id, "bc")
    bot.send_message(msg.chat.id, "📢 Xabar matnini yozing:")

@bot.message_handler(commands=["addbalance"])
def cmd_addbal(msg):
    if msg.from_user.id != ADMIN_ID: return
    try:
        p = msg.text.split(); tid = int(p[1]); amt = int(p[2])
        add_bal(tid, amt)
        bot.send_message(msg.chat.id, f"✅ {tid} ga {amt:,} so'm qo'shildi!")
        bot.send_message(tid, f"💰 Hisobingizga {amt:,} so'm qo'shildi!\nBalans: {get_balance(tid):,} so'm")
    except: bot.send_message(msg.chat.id, "❌ /addbalance [id] [summa]")

@bot.message_handler(commands=["done"])
def cmd_done(msg):
    uid = msg.from_user.id
    imgs = UI.get(uid, [])
    if not imgs: return
    pm = bot.send_message(uid, "⏳ PDF yaratilmoqda...")
    td = tempfile.mkdtemp()
    try:
        out = os.path.join(td, "r.pdf")
        if imgs_to_pdf(imgs, out):
            with open(out, "rb") as f:
                bot.send_document(uid, f, caption="📄 PDF fayl!")
            log_act(uid, "conv", "img_pdf")
        else: bot.send_message(uid, "❌ Xatolik.")
    finally:
        shutil.rmtree(td, ignore_errors=True)
        UI.pop(uid, None); cst(uid)
    try: bot.delete_message(uid, pm.message_id)
    except: pass
    bot.send_message(uid, "✅", reply_markup=main_kb(uid))

# ============================================================
# RASM HANDLER
# ============================================================
@bot.message_handler(content_types=["photo"])
def photo_h(msg):
    uid = msg.from_user.id
    if not check_subscription(uid) and get_sub_channels():
        bot.send_message(uid, "⚠️ *Avval kanallarga obuna bo\'ling!*",
            parse_mode="Markdown", reply_markup=sub_check_kb())
        return
    state = gst(uid)
    ud = UD.get(uid, {})

    # To'lov cheki qabul qilish
    if state == "wait_topup_receipt":
        req_id = ud.get("topup_req_id")
        amount = ud.get("topup_amount", 0)
        user = get_user(uid)
        fname = user["first_name"] if user else "Noma'lum"
        uname = user["username"] if user and user.get("username") else "username_yoq"

        # Adminга chek + ma'lumotlar
        admin_kb = types.InlineKeyboardMarkup(row_width=2)
        admin_kb.add(
            types.InlineKeyboardButton(
                "✅ Tasdiqlash", callback_data=f"topup_ok:{req_id}:{uid}:{amount}"),
            types.InlineKeyboardButton(
                "❌ Rad etish", callback_data=f"topup_no:{req_id}:{uid}")
        )
        caption = (
            f"💳 *Yangi to'ldirish so'rovi* #{req_id}\n\n"
            f"👤 {fname}\n"
            f"📱 @{uname}\n"
            f"🆔 `{uid}`\n"
            f"💰 Summa: *{amount:,} so'm*"
        )
        try:
            bot.send_photo(ADMIN_ID,
                msg.photo[-1].file_id,
                caption=caption,
                parse_mode="Markdown",
                reply_markup=admin_kb)
        except Exception as e:
            logger.error(f"Admin receipt: {e}")
            bot.send_message(ADMIN_ID, caption,
                parse_mode="Markdown", reply_markup=admin_kb)

        cst(uid)
        bot.send_message(uid,
            f"✅ Chekingiz qabul qilindi!\n\n"
            f"💰 Summa: *{amount:,} so'm*\n"
            f"⏳ So'rovingiz ko'rib chiqilmoqda...\n"
            f"Tez orada hisobingizga tushadi!",
            parse_mode="Markdown", reply_markup=main_kb(uid))
        return
    state = gst(uid)
    ph = msg.photo[-1]
    td = tempfile.mkdtemp()
    try:
        fi = bot.get_file(ph.file_id)
        data = bot.download_file(fi.file_path)
        p = os.path.join(td, f"img_{uid}_{len(UI.get(uid,[]))}.jpg")
        with open(p, "wb") as f: f.write(data)
        if state == "img":
            UI.setdefault(uid, []).append(p)
            bot.send_message(uid, t(uid, "img_count", n=len(UI[uid])))
        elif state and "wait_img" in state:
            UI.setdefault(uid, []).append(p)
            n = len(UI[uid]) - 1
            pkb = types.InlineKeyboardMarkup(row_width=5)
            btns = [types.InlineKeyboardButton(str(i), callback_data=f"ipage:{n}:{i}") for i in range(1,16)]
            pkb.add(*btns)
            pkb.add(types.InlineKeyboardButton("✅ Davom etish", callback_data="img_done"))
            bot.send_message(uid, t(uid, "img_accept"), reply_markup=pkb)
        else:
            bot.send_message(uid, "❌ Hozir rasm kutilmayapti.")
    except Exception as e:
        logger.error(f"Photo: {e}")

# ============================================================
# FAYL HANDLER
# ============================================================
@bot.message_handler(content_types=["document"])
def doc_h(msg):
    uid = msg.from_user.id
    state = gst(uid)
    d = msg.document
    if not d: return
        
    # Kitob PDF yuklash
    if state == "wait_book_pdf":
        pm_b = bot.send_message(uid, "⏳ PDF o'qilmoqda...")
        try:
            fi_b = bot.get_file(d.file_id)
            data_b = bot.download_file(fi_b.file_path)
            inp_b = os.path.join(td, d.file_name or "book.pdf")
            with open(inp_b, "wb") as f: f.write(data_b)
            txt_b = pdf_to_text(inp_b) if (d.file_name or "").lower().endswith(".pdf") else ""
            book_name_b = (d.file_name or "kitob").replace(".pdf","").replace(".PDF","")
            if txt_b and len(txt_b) > 100:
                UD.setdefault(uid, {})["source_text"] = txt_b[:4000]
                UD[uid]["source_type"] = "pdf"
            else:
                UD.setdefault(uid, {})["source_type"] = "text"
            UD.setdefault(uid, {})["book_name"] = book_name_b
            svc_b = ud.get("svc", "referat")
            sst(uid, f"{svc_b}_t")
            try: bot.delete_message(uid, pm_b.message_id)
            except: pass
            bot.send_message(uid,
                f"✅ Kitob: *{book_name_b}*\n\n📝 Endi mavzuni kiriting:",
                parse_mode="Markdown", reply_markup=bk_kb())
        except Exception as e_b:
            logger.error(f"Book PDF: {e_b}")
            bot.send_message(uid, "❌ Xato. Kitob nomini yozing:")
            sst(uid, "wait_book_name")
        shutil.rmtree(td, ignore_errors=True)
        return

    if d.file_size > 20*1024*1024:
        bot.send_message(uid, "❌ Fayl juda katta (max 20MB)"); return
    fname = (d.file_name or "").lower()
    td = tempfile.mkdtemp()
    try:
        fi = bot.get_file(d.file_id)
        data = bot.download_file(fi.file_path)
        inp = os.path.join(td, d.file_name or "f")
        with open(inp, "wb") as f: f.write(data)

        if state == "imlo_f":
            pm = bot.send_message(uid, t(uid, "imlo_check"))
            txt = ""
            if fname.endswith(".txt"):
                with open(inp, "r", encoding="utf-8", errors="ignore") as f: txt = f.read()
            elif fname.endswith(".pdf"):
                txt = pdf_to_text(inp)
            if txt:
                fixed = fix_spell(txt, get_lang(uid))
                try: bot.delete_message(uid, pm.message_id)
                except: pass
                bot.send_message(uid, t(uid, "imlo_done") + fixed[:3500], parse_mode="Markdown")
            else:
                try: bot.delete_message(uid, pm.message_id)
                except: pass
                bot.send_message(uid, t(uid, "imlo_notfound"))
            cst(uid)
            bot.send_message(uid, t(uid, "menu"), reply_markup=main_kb(uid))

        elif state == "cv_pptx":
            pm = bot.send_message(uid, t(uid, "converting"))
            if fname.endswith(".pptx"):
                try:
                    from pptx import Presentation as _P
                    prs = _P(inp)
                    content_txt = "\n".join(
                        shape.text for slide in prs.slides
                        for shape in slide.shapes if shape.has_text_frame
                    )
                    out, td2 = make_pdf(content_txt, fname, {})
                    if out:
                        with open(out, "rb") as f:
                            bot.send_document(uid, f, caption="📄 PDF tayyor!")
                        shutil.rmtree(td2, ignore_errors=True)
                    else: bot.send_message(uid, "❌ Xatolik.")
                except Exception as e:
                    bot.send_message(uid, f"❌ {e}")
            else: bot.send_message(uid, t(uid, "wrong_format"))
            try: bot.delete_message(uid, pm.message_id)
            except: pass
            cst(uid); bot.send_message(uid, t(uid, "menu"), reply_markup=main_kb(uid))
    except Exception as e:
        logger.error(f"Doc: {e}")
    finally:
        shutil.rmtree(td, ignore_errors=True)

# ============================================================
# MATN HANDLER
# ============================================================
@bot.message_handler(content_types=["text"])
def text_h(msg):
    uid = msg.from_user.id
    text = msg.text.strip()
    state = gst(uid)
    ud = UD.get(uid, {})

    reg_user(uid, msg.from_user.username or "", msg.from_user.first_name or "")

    # Bekor qilish — to'lov jarayonida
    if text == "❌ Bekor qilish" and gst(uid) == "wait_topup_receipt":
        req_id = UD.get(uid, {}).get("topup_req_id")
        if req_id:
            reject_topup(req_id)
        cst(uid)
        bot.send_message(uid, "❌ To'ldirish bekor qilindi.", reply_markup=main_kb(uid))
        return

    # Obuna tekshirish (til tanlash va /start bundan mustasno)
    skip_sub = [t_val for lang in ["uz","ru","en"] for t_val in [
        TEXTS[lang].get("btn_lang_uz",""), TEXTS[lang].get("btn_lang_ru",""), 
        TEXTS[lang].get("btn_lang_en","")
    ]]
    if not check_subscription(uid) and text not in skip_sub and state not in ["lang_select"]:
        channels = get_sub_channels()
        if channels:
            bot.send_message(uid,
                "⚠️ *Botdan foydalanish uchun quyidagi kanallarga obuna bo\'ling!*",
                parse_mode="Markdown", reply_markup=sub_check_kb())
            return

    # Barcha tillardagi menyu tugmalari
    menu_map = {}
    for lang in ["uz","ru","en"]:
        tx = TEXTS[lang]
        menu_map[tx["btn_referat"]] = ("referat_t", "referat")
        menu_map[tx["btn_kurs"]] = ("kurs_t", "kurs")
        menu_map[tx["btn_mustaqil"]] = ("mustaqil_t", "mustaqil")
        menu_map[tx["btn_maqola"]] = ("maqola_t", "maqola")
        menu_map[tx["btn_prez"]] = ("prez_t", "prez")
        menu_map[tx["btn_test"]] = ("test_t", "test")

    # Menyu tugmalari
    if text in menu_map:
        st2, svc = menu_map[text]
        sst(uid, st2, svc=svc)
        # Referat, kurs, mustaqil, maqola uchun manba tanlash
        if svc in ("referat", "kurs", "mustaqil", "maqola"):
            bot.send_message(uid, "📚 Qaysi manbadan foydalanaylik?", reply_markup=source_kb())
        else:
            bot.send_message(uid, t(uid, "enter_topic"), reply_markup=bk_kb())
        return

    # Imlo
    imlo_btns = [TEXTS[l]["btn_imlo"] for l in ["uz","ru","en"]]
    if text in imlo_btns:
        sst(uid, "imlo_t")
        kb2 = types.InlineKeyboardMarkup()
        kb2.add(types.InlineKeyboardButton("📁 Fayl yuborish (PDF/TXT)", callback_data="imlo_file"))
        bot.send_message(uid, t(uid, "imlo_prompt"), reply_markup=kb2); return

    # Konvertatsiya
    konv_btns = [TEXTS[l]["btn_konv"] for l in ["uz","ru","en"]]
    if text in konv_btns:
        bot.send_message(uid, "🔄 Format tanlang:", reply_markup=conv_kb()); return

    # Balans
    bal_btns = [TEXTS[l]["btn_balans"] for l in ["uz","ru","en"]]
    if text in bal_btns:
        bal = get_balance(uid)
        total_ref, paid_ref, pending_ref = get_referral_stats(uid)
        ref_link = get_referral_link(uid)
        txt_bal = (
            f"💳 *Hisobingiz*\n\n"
            f"💰 Balans: *{bal:,} so'm*\n"
            f"👥 Referal: *{paid_ref}* ta (bonus olindi)\n"
            f"⏳ Kutilmoqda: *{pending_ref}* ta\n\n"
            f"🔗 Referal havola:\n`{ref_link}`"
        )
        kb2 = types.InlineKeyboardMarkup(row_width=1)
        kb2.add(
            types.InlineKeyboardButton("💳 Balans to'ldirish", callback_data="topup"),
            types.InlineKeyboardButton("👥 Referal statistika", callback_data="ref:stats"),
            types.InlineKeyboardButton("🔗 Referal havolam", callback_data="ref:link")
        )
        bot.send_message(uid, txt_bal, parse_mode="Markdown", reply_markup=kb2); return

    # Referal
    ref_btns = [TEXTS[l].get("btn_referral","") for l in ["uz","ru","en"]]
    if text in ref_btns and text:
        total_ref, paid_ref, pending_ref = get_referral_stats(uid)
        ref_link = get_referral_link(uid)
        earned = paid_ref * REFERAL_BONUS
        txt_ref = (
            f"👥 *Referal tizimi*\n\n"
            f"🎁 Har bir do'st uchun: *{REFERAL_BONUS:,} so'm*\n"
            f"📋 Do'st birinchi buyurtma bergandan keyin bonus beriladi\n\n"
            f"📊 *Statistika:*\n"
            f"👤 Jami taklif qilganlar: *{total_ref}* ta\n"
            f"✅ Bonus olindi: *{paid_ref}* ta\n"
            f"⏳ Kutilmoqda: *{pending_ref}* ta\n"
            f"💰 Jami topilgan: *{earned:,} so'm*\n\n"
            f"🔗 *Sizning havolangiz:*\n`{ref_link}`\n\n"
            f"Bu havolani do'stlaringizga yuboring!"
        )
        bot.send_message(uid, txt_ref, parse_mode="Markdown", reply_markup=referral_kb(uid)); return

    # Buyurtmalarim
    orders_btns = [TEXTS[l]["btn_orders"] for l in ["uz","ru","en"]]
    if text in orders_btns:
        rows = get_buyurtmalar(uid)
        if not rows:
            bot.send_message(uid, t(uid, "no_orders"), reply_markup=main_kb(uid)); return
        tur_n = {"referat":"📄 Referat","kurs":"📝 Kurs ishi","mustaqil":"📋 Mustaqil ish",
                 "maqola":"📰 Maqola","prez":"📊 Prezentatsiya","test":"✅ Test"}
        txt2 = t(uid, "orders_title")
        kb_o = types.InlineKeyboardMarkup(row_width=1)
        has_pending = False
        for row in rows:
            # row: id, tur, mavzu, fmt, sah, narx, status, sana
            if len(row) == 8:
                order_id, tur, mavzu, fmt, sah, narx, status, sana = row
            else:
                # Eski format
                tur, mavzu, fmt, sah, narx, sana = row[:6]
                order_id, status = None, "done"
            tur_lbl = tur_n.get(tur, tur)
            sl = t(uid,"slide_pages") if tur=="prez" else (t(uid,"savol_pages") if tur=="test" else t(uid,"bet_pages"))
            status_icon = "⏳" if status == "pending" else "✅"
            txt2 += f"{status_icon} {tur_lbl}\n📌 {mavzu}\n📁 {fmt.upper()} | {sah} {sl} | 💰 {int(narx):,} so'm\n🕐 {sana}\n\n"
            if status == "pending" and order_id:
                has_pending = True
                kb_o.add(types.InlineKeyboardButton(
                    f"▶️ Davom ettirish: {mavzu[:25]}", callback_data=f"resume:{order_id}"))
        if not has_pending:
            kb_o = main_kb(uid)
        bot.send_message(uid, txt2, parse_mode="Markdown", reply_markup=kb_o); return

    # Donat
    donat_btns = [TEXTS[l]["btn_donat"] for l in ["uz","ru","en"]]
    if text in donat_btns:
        kb2 = types.InlineKeyboardMarkup()
        kb2.add(types.InlineKeyboardButton("🌐 Donat", url=DONATE_URL))
        bot.send_message(uid,
            f"💝 *Donat*\n\n💳 Karta: `{DONATE_CARD}`\n🟢 Click: `{DONATE_CLICK}`",
            parse_mode="Markdown", reply_markup=kb2); return

    # Yordam
    help_btns = [TEXTS[l]["btn_help"] for l in ["uz","ru","en"]]
    if text in help_btns:
        help_txt = t(uid, "help_text") + "\n\n" + t(uid, "prices",
            p1=PRICE_PAGE, p2=PRICE_KURS, p3=PRICE_MUSTAQIL,
            p4=PRICE_MAQOLA, p5=PRICE_SLIDE, p6=PRICE_TEST)
        bot.send_message(uid, help_txt, parse_mode="Markdown", reply_markup=main_kb(uid)); return

    # Admin
    admin_btns = [TEXTS[l]["btn_admin"] for l in ["uz","ru","en"]]
    if text in admin_btns:
        kb2 = types.InlineKeyboardMarkup()
        kb2.add(types.InlineKeyboardButton("💬 Adminga yozish",
            url=f"https://t.me/{ADMIN_USERNAME.lstrip('@')}"))
        bot.send_message(uid, "👨‍💼 Admin", reply_markup=kb2); return

    # Broadcast


    # Buyurtma mavzusini tahrirlash
    if state == "edit_order_topic":
        order_id = ud.get("edit_order_id")
        if order_id:
            try:
                c = sqlite3.connect("edubot.db"); cur = c.cursor()
                cur.execute("UPDATE buyurtmalar SET mavzu=? WHERE id=?", (text, order_id))
                c.commit(); c.close()
                bot.send_message(uid, f"✅ Mavzu yangilandi: *{text}*\n\nDavom etish uchun Buyurtmalarim bo\'limiga o\'ting.",
                    parse_mode="Markdown", reply_markup=main_kb(uid))
            except: bot.send_message(uid, "❌ Xato.")
        cst(uid)
        return

    # Buyurtma bet/slayd sonini tahrirlash
    if state == "edit_order_pages":
        order_id = ud.get("edit_order_id")
        if order_id and text.isdigit():
            pages = int(text)
            try:
                c = sqlite3.connect("edubot.db"); cur = c.cursor()
                cur.execute("SELECT tur,narx FROM buyurtmalar WHERE id=?", (order_id,))
                row = cur.fetchone(); c.close()
                if row:
                    tur, old_narx = row
                    price = PRICE_SLIDE if tur == "prez" else PRICE_PAGE
                    new_narx = pages * price
                    c2 = sqlite3.connect("edubot.db"); cur2 = c2.cursor()
                    cur2.execute("UPDATE buyurtmalar SET sahifalar=?, narx=? WHERE id=?", (str(pages), new_narx, order_id))
                    c2.commit(); c2.close()
                    sl = "slayd" if tur == "prez" else "bet"
                    bot.send_message(uid, f"✅ {pages} {sl}, narx: {new_narx:,} so'm\n\nBuyurtmalarim bo\'limidan davom eting.",
                        parse_mode="Markdown", reply_markup=main_kb(uid))
            except: bot.send_message(uid, "❌ Xato.")
        else:
            bot.send_message(uid, "❌ Raqam kiriting!", reply_markup=main_kb(uid))
        cst(uid)
        return

    # Kitob nomi kutish
    if state == "wait_book_name":
        UD.setdefault(uid, {})["book_name"] = text
        svc = ud.get("svc", "referat")
        sst(uid, f"{svc}_t")
        bot.send_message(uid, f"✅ Kitob: *{text}*\n\n📝 Endi mavzuni kiriting:", 
            parse_mode="Markdown", reply_markup=bk_kb())
        return

    if state == "bc":
        users = all_users()
        ok = 0
        for u2 in users:
            try: bot.send_message(u2, text); ok += 1
            except: pass
        cst(uid)
        bot.send_message(uid, f"✅ {ok}/{len(users)} ta yuborildi!"); return

    # INFO STATES
    if state in INFO_STATES:
        idx = INFO_STATES.index(state)
        _, field, required = INFO_STEPS[idx]
        val = text.strip()
        if val: UD.setdefault(uid, {})[field] = val
        if idx < len(INFO_STEPS) - 1:
            ns, nfield, nreq = INFO_STEPS[idx + 1]
            sst(uid, ns)
            msg_txt = t(uid, ns)
            if not nreq:
                msg_txt += f"\n{t(uid,'optional')}"
                bot.send_message(uid, msg_txt, reply_markup=skip_kb(ns))
            else:
                bot.send_message(uid, msg_txt, reply_markup=bk_kb())
        else:
            finish_info(uid, UD.get(uid, {}))
        return

    # MAVZU KIRITISH STATES
    svc_topic_states = ["referat_t","kurs_t","mustaqil_t","maqola_t","prez_t","test_t"]
    if state in svc_topic_states:
        svc = ud.get("svc", state.replace("_t",""))
        UD.setdefault(uid, {})["topic"] = text
        # Ma'lumotlarni so'rash
        sst(uid, "ask_name", svc=svc, topic=text)
        bot.send_message(uid, t(uid, "ask_name"), reply_markup=bk_kb()); return

    # Sahifa soni
    if state in ["referat_p","kurs_p","mustaqil_p","maqola_p"]:
        try:
            pages = int(text)
            if pages < 1 or pages > 100:
                bot.send_message(uid, "❌ 1-100 oralig'ida kiriting!"); return
            svc = ud.get("svc","referat")
            price_map = {"referat":PRICE_PAGE,"kurs":PRICE_KURS,"mustaqil":PRICE_MUSTAQIL,"maqola":PRICE_MAQOLA}
            price = price_map.get(svc, PRICE_PAGE)
            total = pages * price
            UD.setdefault(uid, {})["pages"] = pages
            UD[uid]["total"] = total
            bal = get_balance(uid)
            if bal < total:
                save_pending_and_notify(uid, svc, UD.get(uid,{}).get("topic",""), "docx", pages, total, UD.get(uid,{}))
                return
            sst(uid, f"{svc}_lang")
            bot.send_message(uid,
                f"✅ {pages} bet × {price:,} = *{total:,} so'm*\n\n{t(uid,'ask_lang')}",
                parse_mode="Markdown", reply_markup=lc_kb(f"{svc}_lang"))
        except:
            bot.send_message(uid, "❌ Raqam kiriting!")
        return

    # Test savol soni
    if state == "test_p":
        try:
            count = int(text)
            if count < 1 or count > 1000:
                bot.send_message(uid, "❌ 1-1000 oralig'ida kiriting!"); return
            total = count * PRICE_TEST
            UD.setdefault(uid, {})["count"] = count
            UD[uid]["total"] = total
            bal = get_balance(uid)
            if bal < total:
                save_pending_and_notify(uid, "test", UD.get(uid,{}).get("topic",""), "txt", count, total, UD.get(uid,{}))
                return
            sst(uid, "test_lang")
            bot.send_message(uid,
                f"✅ {count} savol × {PRICE_TEST:,} = *{total:,} so'm*\n\n{t(uid,'ask_lang')}",
                parse_mode="Markdown", reply_markup=lc_kb("test_lang"))
        except:
            bot.send_message(uid, "❌ Raqam kiriting!")
        return

    # Slayd soni (custom)
    if state == "prez_slides_custom":
        try:
            slides = int(text)
            if slides < 5 or slides > 100:
                bot.send_message(uid, "❌ 5-100 oralig'ida kiriting!"); return
            total = slides * PRICE_SLIDE
            UD.setdefault(uid, {})["slides"] = slides
            UD[uid]["total"] = total
            bal = get_balance(uid)
            if bal < total:
                save_pending_and_notify(uid, "prez", UD.get(uid,{}).get("topic",""), "pptx", slides, total, UD.get(uid,{}))
                return
            sst(uid, "prez_plans")
            bot.send_message(uid,
                f"✅ {slides} slayd × {PRICE_SLIDE:,} = *{total:,} so'm*\n\n{t(uid,'ask_plans')}",
                parse_mode="Markdown", reply_markup=plans_kb())
        except:
            bot.send_message(uid, "❌ Raqam kiriting!")
        return

    # Reja soni (custom)
    if state == "prez_plans_custom":
        try:
            plans = int(text)
            if plans < 2 or plans > 20:
                bot.send_message(uid, "❌ 2-20 oralig'ida kiriting!"); return
            UD.setdefault(uid, {})["plans_count"] = plans
            sst(uid, "prez_lang")
            bot.send_message(uid, t(uid, "ask_lang"), reply_markup=lc_kb("prez_lang"))
        except:
            bot.send_message(uid, "❌ Raqam kiriting!")
        return

    # Imlo matn
    if state == "imlo_t":
        pm = bot.send_message(uid, t(uid, "imlo_check"))
        fixed = fix_spell(text, get_lang(uid))
        try: bot.delete_message(uid, pm.message_id)
        except: pass
        bot.send_message(uid, t(uid, "imlo_done") + fixed[:3500], parse_mode="Markdown")
        cst(uid)
        bot.send_message(uid, t(uid, "menu"), reply_markup=main_kb(uid)); return

# ============================================================
# CALLBACK HANDLER
# ============================================================
@bot.callback_query_handler(func=lambda c: True)
def cb(call):
    uid = call.from_user.id
    d = call.data
    ud = UD.get(uid, {})

    # Obuna tekshirish (check_sub va lang bundan mustasno)
    sub_exempt = ["check_sub", "bk"] + [f"lang:{l}" for l in ["uz","ru","en"]]
    if d not in sub_exempt and not d.startswith("lang:") and not check_subscription(uid):
        channels = get_sub_channels()
        if channels:
            try: bot.answer_callback_query(call.id, "⚠️ Avval kanallarga obuna bo\'ling!")
            except: pass
            bot.send_message(uid,
                "⚠️ *Botdan foydalanish uchun quyidagi kanallarga obuna bo\'ling!*",
                parse_mode="Markdown", reply_markup=sub_check_kb())
            return

    # Til tanlash
    if d.startswith("lang:"):
        try: bot.delete_message(uid, call.message.message_id)
        except: pass
        lang = d[5:]
        set_lang(uid, lang)
        try: bot.edit_message_text(TEXTS[lang]["lang_set"], uid, call.message.message_id)
        except: pass
        bot.send_message(uid, t(uid, "menu"), reply_markup=main_kb(uid))
        return

    # Noop
    if d == "noop":
        try: bot.answer_callback_query(call.id)
        except: pass
        return

    # Manba tanlash (referat/kurs/mustaqil/maqola)
    if d.startswith("src:"):
        try: bot.delete_message(uid, call.message.message_id)
        except: pass
        src = d[4:]
        UD.setdefault(uid, {})["source_type"] = src
        svc = ud.get("svc", "referat")
        if src == "text":
            sst(uid, "wait_book_name")
            bot.send_message(uid,
                "📚 Kitob yoki fan nomini yozing:\n"
                "(Masalan: 'Biologiya 10-sinf' yoki 'Iqtisodiyot asoslari')",
                reply_markup=bk_kb())
        elif src == "pdf":
            sst(uid, "wait_book_pdf")
            bot.send_message(uid, "📄 PDF kitobni yuboring:", reply_markup=bk_kb())
        else:
            # Manbasisiz — to'g'ri mavzuga o'tish
            UD[uid]["source_type"] = "none"
            sst(uid, f"{svc}_t", svc=svc)
            bot.send_message(uid, t(uid, "enter_topic"), reply_markup=bk_kb())
        return

    # Diagramma tanlash
    if d.startswith("diag:"):
        try: bot.delete_message(uid, call.message.message_id)
        except: pass
        UD.setdefault(uid, {})["with_diagram"] = (d == "diag:yes")
        sst(uid, "prez_img")
        bot.send_message(uid, "🖼 Prezentatsiyaga rasm qo'shmoqchimisiz?", reply_markup=img_choice_kb())
        return

    # Menyu
    if d == "bk":
        cst(uid)
        try: bot.edit_message_text(t(uid,"menu"), uid, call.message.message_id)
        except: pass
        bot.send_message(uid, t(uid,"menu"), reply_markup=main_kb(uid)); return

    # Orqaga
    if d == "back_step":
        prev = go_back(uid)
        if prev:
            bot.send_message(uid, f"◀️ {prev}", reply_markup=bk_kb())
        else:
            cst(uid)
            bot.send_message(uid, t(uid,"menu"), reply_markup=main_kb(uid)); return

    # Skip
    if d.startswith("skip:"):
        try: bot.delete_message(uid, call.message.message_id)
        except: pass
        ns = d[5:]
        idx = INFO_STATES.index(ns) if ns in INFO_STATES else -1
        if idx >= 0 and idx < len(INFO_STEPS) - 1:
            next_ns, _, nreq = INFO_STEPS[idx + 1]
            sst(uid, next_ns)
            msg_txt = t(uid, next_ns)
            if not nreq:
                msg_txt += f"\n{t(uid,'optional')}"
                bot.send_message(uid, msg_txt, reply_markup=skip_kb(next_ns))
            else:
                bot.send_message(uid, msg_txt, reply_markup=bk_kb())
        else:
            finish_info(uid, UD.get(uid, {}))
        return

    # Shablon tanlash
    if d.startswith("tmpl_p:"):
        try: bot.delete_message(uid, call.message.message_id)
        except: pass
        page = int(d[7:])
        bot.send_message(uid, "🎨 Shablon tanlang:", reply_markup=tmpl_kb(page))
        return

    if d.startswith("tmpl:"):
        try: bot.delete_message(uid, call.message.message_id)
        except: pass
        tmpl_id = d[5:]
        UD.setdefault(uid, {})["template_id"] = tmpl_id
        sst(uid, "prez_slides")
        bot.send_message(uid,
            t(uid,"enter_slides", price=PRICE_SLIDE),
            reply_markup=slides_kb())
        return

    if d.startswith("tmpl_p:"):
        page = int(d[7:])
        try: bot.edit_message_reply_markup(uid, call.message.message_id, reply_markup=tmpl_kb(page))
        except: pass; return

    # Slayd soni
    if d.startswith("slides:"):
        try: bot.delete_message(uid, call.message.message_id)
        except: pass
        val = d[7:]
        if val == "custom":
            sst(uid, "prez_slides_custom")
            bot.send_message(uid, "✏️ Slayd sonini kiriting (5-100):"); return
        slides = int(val)
        total = slides * PRICE_SLIDE
        UD.setdefault(uid, {})["slides"] = slides
        UD[uid]["total"] = total
        bal = get_balance(uid)
        if bal < total:
            save_pending_and_notify(uid, "prez", ud.get("topic",""), "pptx", slides, total, ud)
            return
        sst(uid, "prez_plans")
        bot.send_message(uid,
            f"✅ {slides} slayd × {PRICE_SLIDE:,} = *{total:,} so'm*\n\n{t(uid,'ask_plans')}",
            parse_mode="Markdown", reply_markup=plans_kb())
        return

    # Reja soni
    if d.startswith("plans:"):
        try: bot.delete_message(uid, call.message.message_id)
        except: pass
        plans_count = int(d[6:])
        UD.setdefault(uid, {})["plans_count"] = plans_count
        # Diagramma savolini ko'rsatish
        sst(uid, "prez_diag")
        diag_kb = types.InlineKeyboardMarkup(row_width=1)
        diag_kb.add(
            types.InlineKeyboardButton("📊 Ha, diagramma qo'shilsin", callback_data="diag:yes"),
            types.InlineKeyboardButton("📝 Yo'q, faqat matn", callback_data="diag:no")
        )
        bot.send_message(uid,
            f"✅ {plans_count} ta bo'lim\n\n"
            f"📊 Diagramma va infografika qo'shilsinmi?",
            parse_mode="Markdown", reply_markup=diag_kb)
        return
        val = d[6:]
        if val == "custom":
            sst(uid, "prez_plans_custom")
            bot.send_message(uid, "✏️ Bo'lim sonini kiriting (2-20):"); return
        UD.setdefault(uid, {})["plans_count"] = int(val)
        sst(uid, "prez_lang")
        bot.send_message(uid, t(uid,"ask_lang"), reply_markup=lc_kb("prez_lang"))
        return

    # Til (til_lang format: referat_lang:uz)
    for svc in ["referat","kurs","mustaqil","maqola","prez","test"]:
        if d.startswith(f"{svc}_lang:"):
            lang = d.split(":")[1]
            UD.setdefault(uid, {})["lang"] = lang
            if svc == "prez":
                sst(uid, "prez_diag")
                diag_kb = types.InlineKeyboardMarkup(row_width=1)
                diag_kb.add(
                    types.InlineKeyboardButton("📊 Ha, diagramma qo'shilsin", callback_data="diag:yes"),
                    types.InlineKeyboardButton("📝 Yo'q, faqat matn", callback_data="diag:no")
                )
                bot.send_message(uid,
                    "📊 *Diagramma va infografika qo'shilsinmi?*\n\n"
                    "Agar mavzuda raqamli statistik ma'lumotlar bo'lsa diagramma qo'shiladi\n"
                    "(iqtisod, biologiya, tibbiyot, fizika va h.k.)",
                    parse_mode="Markdown", reply_markup=diag_kb)
            elif svc == "test":
                sst(uid, "test_confirm")
                count = ud.get("count", 10)
                total = ud.get("total", count * PRICE_TEST)
                bot.send_message(uid,
                    f"📊 *Test:* {ud.get('topic','')}\n"
                    f"📝 {count} ta savol | 💰 {total:,} so'm\n\n"
                    f"Tasdiqlaysizmi?",
                    parse_mode="Markdown",
                    reply_markup=types.InlineKeyboardMarkup().add(
                        types.InlineKeyboardButton("✅ Ha, yaratish", callback_data="test_go"),
                        types.InlineKeyboardButton("❌ Bekor", callback_data="bk")
                    ))
            else:
                sst(uid, f"{svc}_pages")
                price_map = {"referat":PRICE_PAGE,"kurs":PRICE_KURS,"mustaqil":PRICE_MUSTAQIL,"maqola":PRICE_MAQOLA}
                price = price_map.get(svc, PRICE_PAGE)
                bot.send_message(uid, t(uid,"enter_pages",price=price), reply_markup=bk_kb())
            return

    # Test savol soni
    if d.startswith("tcount:"):
        try: bot.delete_message(uid, call.message.message_id)
        except: pass
        val = d[7:]
        if val == "custom":
            sst(uid, "test_p")
            bot.send_message(uid, t(uid,"enter_count",price=PRICE_TEST)); return
        count = int(val)
        total = count * PRICE_TEST
        UD.setdefault(uid, {})["count"] = count
        UD[uid]["total"] = total
        bal = get_balance(uid)
        if bal < total:
            save_pending_and_notify(uid, "test", ud.get("topic",""), "txt", count, total, ud)
            return
        sst(uid, "test_lang")
        bot.send_message(uid,
            f"✅ {count} savol × {PRICE_TEST:,} = *{total:,} so'm*\n\n{t(uid,'ask_lang')}",
            parse_mode="Markdown", reply_markup=lc_kb("test_lang"))
        return

    # Rasm tanlash
    if d.startswith("img:"):
        try: bot.delete_message(uid, call.message.message_id)
        except: pass
        choice = d[4:]
        slides = ud.get("slides", 15)
        if choice == "ai":
            ai_slides = list(range(3, slides+1, 3))
            UD.setdefault(uid, {})["ai_img_slides"] = ai_slides
            sst(uid, "prez_fmt")
            bot.send_message(uid, t(uid,"ask_format"), reply_markup=prez_fmt_kb())
        elif choice == "user":
            sst(uid, "wait_img")
            bot.send_message(uid,
                f"🖼 Rasmlarni yuboring (har rasm uchun qaysi slayd belgilanadi)\n"
                f"Tugatgach /done yozing.")
        elif choice == "none":
            UD.setdefault(uid, {})["ai_img_slides"] = []
            sst(uid, "prez_fmt")
            bot.send_message(uid, t(uid,"ask_format"), reply_markup=prez_fmt_kb())
        return

    if d == "img_done":
        sst(uid, "prez_fmt")
        bot.send_message(uid, t(uid,"ask_format"), reply_markup=prez_fmt_kb()); return

    if d.startswith("ipage:"):
        parts = d.split(":"); img_idx = int(parts[1]); page_num = int(parts[2])
        UD.setdefault(uid, {}).setdefault("img_pages", {})[str(img_idx)] = page_num
        bot.answer_callback_query(call.id, f"✅ {page_num}-slaydga qo'shildi!")
        return

    # Format tanlash (hujjat)
    for svc in ["referat","kurs","mustaqil","maqola"]:
        if d.startswith(f"{svc}_fmt:"):
            fmt = d.split(":")[1]
            UD.setdefault(uid, {})["fmt"] = fmt
            topic = ud.get("topic","")
            pages = ud.get("pages", 5)
            lang = ud.get("lang", get_lang(uid))
            total = ud.get("total", pages * PRICE_PAGE)
            bal = get_balance(uid)
            if bal < total:
                save_pending_and_notify(uid, svc, topic, fmt, pages, total, ud)
                return
            deduct(uid, total)
            pm = bot.send_message(uid, t(uid,"preparing"))
            import threading
            def gen_doc_task():
                td2 = None
                try:
                    content = gen_doc(svc, topic, pages, lang, ud)
                    if fmt == "docx":
                        out, td2 = make_docx(content, topic, ud)
                        cap = "📝 DOCX fayl tayyor!"
                        fname2 = "dokument.docx"
                    elif fmt == "pdf":
                        out, td2 = make_pdf(content, topic, ud)
                        cap = "📄 PDF fayl tayyor!"
                        fname2 = "dokument.pdf"
                    else:
                        out = None
                        bot.send_message(uid, f"📱 *{topic}*\n\n{content[:4000]}", parse_mode="Markdown")
                        cap = None
                    if out and cap:
                        with open(out, "rb") as f:
                            bot.send_document(uid, f, caption=cap, visible_file_name=fname2)
                    log_act(uid, svc, topic, total)
                    save_buyurtma(uid, svc, topic, fmt, pages, total)
                except Exception as e:
                    logger.error(f"Doc gen: {e}")
                    add_bal(uid, total)
                    bot.send_message(uid, t(uid,"error"))
                finally:
                    if td2: shutil.rmtree(td2, ignore_errors=True)
                    try: bot.delete_message(uid, pm.message_id)
                    except: pass
                    cst(uid)
                    bot.send_message(uid, t(uid,"ready",bal=get_balance(uid)), reply_markup=main_kb(uid))
            threading.Thread(target=gen_doc_task).start()
            return

    # Prezentatsiya format
    if d.startswith("pfmt:"):
        try: bot.delete_message(uid, call.message.message_id)
        except: pass
        fmt = d[5:]
        UD.setdefault(uid, {})["fmt"] = fmt
        topic = ud.get("topic","")
        slides = ud.get("slides", 15)
        lang = ud.get("lang", get_lang(uid))
        plans = ud.get("plans_count", 5)
        total = ud.get("total", slides * PRICE_SLIDE)
        tmpl_id = ud.get("template_id", "1")
        tmpl_name = TEMPLATES.get(str(tmpl_id), TEMPLATES["1"])["name"]
        bal = get_balance(uid)
        if bal < total:
            save_pending_and_notify(uid, "prez", topic, fmt, slides, total, ud)
            return
        deduct(uid, total)
        pm = bot.send_message(uid, t(uid,"preparing"))
        user_imgs = UI.get(uid, [])
        img_pages = ud.get("img_pages", {})
        import threading
        def gen_prez_task():
            td2 = None
            try:
                content = gen_prez(topic, slides, lang, ud, plans)
                if fmt in ["pptx","both"]:
                    out, td2 = make_pptx(content, topic, tmpl_id, ud, user_imgs, img_pages)
                    with open(out, "rb") as f:
                        bot.send_document(uid, f, caption="📊 PPTX tayyor!", visible_file_name="prezentatsiya.pptx")
                    shutil.rmtree(td2, ignore_errors=True)
                if fmt in ["html","both"]:
                    out2, td3 = make_html(content, topic, tmpl_id, ud)
                    with open(out2, "rb") as f:
                        bot.send_document(uid, f, caption="🌐 HTML tayyor!", visible_file_name="prezentatsiya.html")
                    shutil.rmtree(td3, ignore_errors=True)
                log_act(uid, "prez", topic, total)
                save_buyurtma(uid, "prez", topic, fmt, slides, total)
            except Exception as e:
                logger.error(f"Prez gen: {e}")
                add_bal(uid, total)
                bot.send_message(uid, t(uid,"error"))
            finally:
                try: bot.delete_message(uid, pm.message_id)
                except: pass
                UI.pop(uid, None); cst(uid)
                bot.send_message(uid, t(uid,"ready",bal=get_balance(uid)), reply_markup=main_kb(uid))
        threading.Thread(target=gen_prez_task).start()
        return

    # Test yaratish
    if d == "test_go":
        try: bot.delete_message(uid, call.message.message_id)
        except: pass
        topic = ud.get("topic","")
        count = ud.get("count", 10)
        lang = ud.get("lang", get_lang(uid))
        total = ud.get("total", count * PRICE_TEST)
        bal = get_balance(uid)
        if bal < total:
            save_pending_and_notify(uid, "test", ud.get("topic",""), "txt", ud.get("count",10), total, ud)
            return
        deduct(uid, total)
        pm = bot.send_message(uid, t(uid,"preparing"))
        import threading
        def gen_test_task():
            try:
                content = gen_test(topic, count, lang)
                bot.send_message(uid, f"✅ *{topic}*\n\n{content[:4000]}", parse_mode="Markdown")
                log_act(uid, "test", topic, total)
                save_buyurtma(uid, "test", topic, "txt", count, total)
            except Exception as e:
                logger.error(f"Test gen: {e}")
                add_bal(uid, total)
                bot.send_message(uid, t(uid,"error"))
            finally:
                try: bot.delete_message(uid, pm.message_id)
                except: pass
                cst(uid)
                bot.send_message(uid, t(uid,"ready",bal=get_balance(uid)), reply_markup=main_kb(uid))
        threading.Thread(target=gen_test_task).start()
        return

    # Konvertatsiya
    if d.startswith("cv:"):
        cv_type = d[3:]
        if cv_type == "img":
            sst(uid, "img")
            bot.send_message(uid, "📷 Rasmlarni yuboring, so'ng /done yozing.")
        elif cv_type == "pptx":
            sst(uid, "cv_pptx")
            bot.send_message(uid, "📊 PPTX faylni yuboring:")
        elif cv_type == "pdf":
            bot.send_message(uid, "🔄 Hozircha PDF → PPTX qo'llab-quvvatlanmaydi.")
        return

    # Imlo fayl
    if d == "imlo_file":
        sst(uid, "imlo_f")
        bot.send_message(uid, "📁 PDF yoki TXT faylni yuboring:"); return

    # Topup
    if d == "topup":
        try: bot.delete_message(uid, call.message.message_id)
        except: pass
        bot.send_message(uid,
            "💳 *Balans to'ldirish*\n\nQuyidagi summalardan birini tanlang:",
            parse_mode="Markdown", reply_markup=topup_kb())
        return

    # To'ldirish miqdori tanlandi
    if d.startswith("topup_amt:"):
        try: bot.delete_message(uid, call.message.message_id)
        except: pass
        amount = int(d[10:])
        # So'rovni saqlab, foydalanuvchiga to'lov ma'lumotlarini ko'rsatamiz
        req_id = save_topup_request(uid, amount)
        UD.setdefault(uid, {})["topup_req_id"] = req_id
        UD[uid]["topup_amount"] = amount
        sst(uid, "wait_topup_receipt")
        bot.send_message(uid,
            f"💳 *To'lov ma'lumotlari*\n\n"
            f"💰 To'lov summasi: *{amount:,} so'm*\n\n"
            f"💳 Karta raqami:\n`{DONATE_CARD}`\n\n"
            f"📱 Click/Payme:\n`{DONATE_CLICK}`\n\n"
            f"━━━━━━━━━━━━━━━━━━━━\n"
            f"✅ To'lovni amalga oshirib, *chek rasmini* yuboring.\n"
            f"So'rovingiz *1 daqiqa* ichida ko'rib chiqiladi!",
            parse_mode="Markdown",
            reply_markup=types.ReplyKeyboardMarkup(resize_keyboard=True).add(
                types.KeyboardButton("❌ Bekor qilish")
            ))
        return

    # Admin: tasdiqlash
    if d.startswith("topup_ok:"):
        if uid != ADMIN_ID:
            try: bot.answer_callback_query(call.id, "❌ Ruxsat yo'q!")
            except: pass
            return
        parts = d.split(":")
        req_id, target_uid, amount = int(parts[1]), int(parts[2]), int(parts[3])
        user_uid, paid_amount = approve_topup(req_id)
        if user_uid:
            # Xabarni yangilash
            try:
                bot.edit_message_caption(
                    caption=f"✅ *Tasdiqlandi* — {paid_amount:,} so'm → {user_uid}",
                    chat_id=ADMIN_ID,
                    message_id=call.message.message_id,
                    parse_mode="Markdown",
                    reply_markup=types.InlineKeyboardMarkup())
            except:
                try:
                    bot.edit_message_reply_markup(ADMIN_ID, call.message.message_id,
                        reply_markup=types.InlineKeyboardMarkup())
                except: pass
            try: bot.answer_callback_query(call.id, f"✅ {paid_amount:,} so'm tasdiqlandi!")
            except: pass
            # Foydalanuvchiga xabar
            try:
                new_bal = get_balance(user_uid)
                bot.send_message(user_uid,
                    f"✅ *Hisobingiz to'ldirildi!*\n\n"
                    f"💰 Qo'shildi: *{paid_amount:,} so'm*\n"
                    f"💳 Joriy balans: *{new_bal:,} so'm*\n\n"
                    f"Xizmatlarimizdan foydalaning! 🎓",
                    parse_mode="Markdown", reply_markup=main_kb(user_uid))
            except: pass
        else:
            try: bot.answer_callback_query(call.id, "❌ So'rov topilmadi!")
            except: pass
        return

    # Admin: rad etish
    if d.startswith("topup_no:"):
        if uid != ADMIN_ID:
            try: bot.answer_callback_query(call.id, "❌ Ruxsat yo'q!")
            except: pass
            return
        parts = d.split(":")
        req_id, target_uid = int(parts[1]), int(parts[2])
        user_uid = reject_topup(req_id)
        try:
            bot.edit_message_reply_markup(ADMIN_ID, call.message.message_id,
                reply_markup=types.InlineKeyboardMarkup())
            bot.send_message(ADMIN_ID, f"❌ #{req_id} rad etildi.")
        except: pass
        if user_uid:
            try:
                bot.send_message(user_uid,
                    f"❌ Afsuski, #{req_id} to'ldirish so'rovingiz rad etildi.\n"
                    f"Murojaat uchun adminга yozing.",
                    reply_markup=main_kb(user_uid))
            except: pass
        return

    # Referal callbacks
    if d == "ref:link":
        try: bot.answer_callback_query(call.id)
        except: pass
        ref_link = get_referral_link(uid)
        bot.send_message(uid,
            f"🔗 *Sizning referal havolangiz:*\n\n`{ref_link}`\n\n"
            f"Bu havolani do'stlaringizga yuboring.\n"
            f"Har biri birinchi buyurtma bergandan so'ng *{REFERAL_BONUS:,} so'm* bonus olasiz!",
            parse_mode="Markdown")
        return

    if d == "ref:stats":
        try: bot.answer_callback_query(call.id)
        except: pass
        total_ref, paid_ref, pending_ref = get_referral_stats(uid)
        earned = paid_ref * REFERAL_BONUS
        bot.send_message(uid,
            f"📊 *Referal statistika*\n\n"
            f"👤 Jami taklif: *{total_ref}* ta\n"
            f"✅ Bonus olindi: *{paid_ref}* ta\n"
            f"⏳ Kutilmoqda: *{pending_ref}* ta\n"
            f"💰 Jami daromad: *{earned:,} so'm*",
            parse_mode="Markdown", reply_markup=referral_kb(uid))
        return

    # Eski topup xabari (endi yuklanmaydi)
    # if d == "topup_old": ...

    try: bot.answer_callback_query(call.id)
    except: pass

# ============================================================
# ASOSIY
# ============================================================
if __name__ == "__main__":
    init_db()
    try:
        bot.set_my_commands([
            types.BotCommand("start", "Botni ishga tushirish"),
            types.BotCommand("referat", "Referat yozish"),
            types.BotCommand("kursishi", "Kurs ishi yozish"),
            types.BotCommand("mustaqilish", "Mustaqil ish yozish"),
            types.BotCommand("maqola", "Ilmiy maqola yozish"),
            types.BotCommand("prezentatsiya", "Prezentatsiya yaratish"),
            types.BotCommand("test", "Test savollari yaratish"),
            types.BotCommand("imlo", "Imlo tuzatish"),
            types.BotCommand("konvertatsiya", "Fayl konvertatsiya"),
            types.BotCommand("balans", "Balansni ko'rish"),
            types.BotCommand("yordam", "Yordam va narxlar"),
            types.BotCommand("menu", "Asosiy menyuni ochish"),
        ])
        logger.info("Buyruqlar ro'yxatlandi!")
    except Exception as e:
        logger.error(f"Commands: {e}")
    logger.info("EduBot v12 ishga tushdi!")
    bot.infinity_polling(timeout=60, long_polling_timeout=60)